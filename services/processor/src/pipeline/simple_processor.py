#!/usr/bin/env python3
"""
Simple Document Processor
Extracts text from various file types and saves as both .txt and .json
Includes file paths in JSON output and handles images/presentations with OCR
"""

import os
import json
import zipfile
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any
import uuid

def setup_dependencies():
    """Check and install required dependencies."""
    deps = [
        ("fitz", "pymupdf", "PDF processing"),
        ("PIL", "pillow", "Image processing"), 
        ("pytesseract", "pytesseract", "OCR"),
        ("pptx", "python-pptx", "PowerPoint"),
        ("docx", "python-docx", "Word docs"),
        ("openpyxl", "openpyxl", "Excel files")
    ]
    
    missing = []
    for module, package, desc in deps:
        try:
            __import__(module)
            print(f"✅ {desc}")
        except ImportError:
            missing.append((package, desc))
            print(f"❌ {desc} - install with: pip install {package}")
    
    if missing:
        print(f"\nInstall missing packages:")
        print(f"pip install {' '.join([p[0] for p in missing])}")
        return False
    return True

class SimpleDocumentProcessor:
    """Simple processor for various document types."""
    
    def __init__(self, output_dir: str):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        
        # Create subdirectories
        self.txt_dir = self.output_dir / "txt_files"
        self.json_dir = self.output_dir / "json_files"
        self.txt_dir.mkdir(exist_ok=True)
        self.json_dir.mkdir(exist_ok=True)
        
        self.stats = {
            "processed": 0,
            "failed": 0,
            "by_type": {}
        }
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """Process a single file and extract text."""
        file_path = Path(file_path)
        ext = file_path.suffix.lower()[1:] if file_path.suffix else "no_extension"
        
        result = {
            "file_path": str(file_path.absolute()),
            "file_name": file_path.name,
            "file_extension": ext,
            "file_size": file_path.stat().st_size if file_path.exists() else 0,
            "processed_at": datetime.now().isoformat(),
            "extracted_text": "",
            "extraction_method": "none",
            "success": False,
            "errors": []
        }
        
        try:
            if ext == "pdf":
                result["extracted_text"], result["extraction_method"] = self._extract_pdf(file_path)
            elif ext in ["png", "jpg", "jpeg", "tiff", "bmp", "gif", "webp"]:
                result["extracted_text"], result["extraction_method"] = self._extract_image(file_path)
            elif ext in ["pptx", "ppt"]:
                result["extracted_text"], result["extraction_method"] = self._extract_powerpoint(file_path)
            elif ext in ["docx", "doc"]:
                result["extracted_text"], result["extraction_method"] = self._extract_word(file_path)
            elif ext in ["xlsx", "xls"]:
                result["extracted_text"], result["extraction_method"] = self._extract_excel(file_path)
            elif ext in ["txt", "md", "csv", "log", "rtf", "html", "htm", "xml", "json"]:
                result["extracted_text"], result["extraction_method"] = self._extract_text(file_path)
            else:
                # Try as text file for unknown extensions
                result["extracted_text"], result["extraction_method"] = self._extract_text(file_path)
                if not result["extracted_text"].strip():
                    result["extracted_text"] = f"[Binary file - could not extract text content]"
                    result["extraction_method"] = "binary_fallback"
                
            result["success"] = len(result["extracted_text"].strip()) > 0
            
        except Exception as e:
            result["errors"].append(f"Processing error: {str(e)}")
            result["extracted_text"] = f"[Error extracting content: {str(e)}]"
            result["extraction_method"] = "error"
            result["success"] = False
        
        return result
    
    def _extract_pdf(self, file_path: Path) -> tuple[str, str]:
        """Extract text from PDF files."""
        try:
            import fitz
            text_parts = []
            
            with fitz.open(str(file_path)) as doc:
                for page_num, page in enumerate(doc):
                    # Extract text
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"=== Page {page_num + 1} ===\n{text}")
                    
                    # Extract images and OCR if page has little text
                    if len(text.strip()) < 100:
                        images = page.get_images()
                        for img_index, img in enumerate(images):
                            try:
                                xref = img[0]
                                pix = fitz.Pixmap(doc, xref)
                                if pix.n - pix.alpha < 4:  # GRAY or RGB
                                    img_data = pix.tobytes("png")
                                    ocr_text = self._ocr_from_bytes(img_data)
                                    if ocr_text.strip():
                                        text_parts.append(f"=== Page {page_num + 1} Image {img_index + 1} (OCR) ===\n{ocr_text}")
                                pix = None
                            except:
                                continue
            
            return "\n\n".join(text_parts), "pymupdf+ocr"
            
        except ImportError:
            return "", "pymupdf_not_available"
        except Exception as e:
            return f"PDF extraction error: {str(e)}", "error"
    
    def _extract_image(self, file_path: Path) -> tuple[str, str]:
        """Extract text from image files using OCR."""
        try:
            from PIL import Image
            import pytesseract
            
            with Image.open(file_path) as img:
                # Preprocess image for better OCR
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # OCR with different PSM modes
                configs = [
                    r'--oem 3 --psm 6',  # Single block
                    r'--oem 3 --psm 11', # Sparse text
                    r'--oem 3 --psm 12'  # Single block, sparse
                ]
                
                best_text = ""
                for config in configs:
                    try:
                        text = pytesseract.image_to_string(img, config=config)
                        if len(text.strip()) > len(best_text.strip()):
                            best_text = text
                    except:
                        continue
                
                return best_text.strip(), "tesseract_ocr"
                
        except ImportError:
            return "", "ocr_not_available"
        except Exception as e:
            return f"Image OCR error: {str(e)}", "error"
    
    def _ocr_from_bytes(self, img_bytes: bytes) -> str:
        """OCR from image bytes."""
        try:
            from PIL import Image
            import pytesseract
            import io
            
            img = Image.open(io.BytesIO(img_bytes))
            text = pytesseract.image_to_string(img, config=r'--oem 3 --psm 11')
            return text.strip()
        except:
            return ""
    
    def _extract_powerpoint(self, file_path: Path) -> tuple[str, str]:
        """Extract text from PowerPoint files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = [f"=== Slide {i + 1} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # Extract table data
                    if shape.has_table:
                        table_text = ["[Table]"]
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_text.append(" | ".join(row_data))
                        slide_content.extend(table_text)
                
                # Add slide notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_content.append(f"[Notes: {notes}]")
                
                slides_text.append("\n".join(slide_content))
            
            return "\n\n".join(slides_text), "python-pptx"
            
        except ImportError:
            return "", "python-pptx_not_available"
        except Exception as e:
            return f"PowerPoint extraction error: {str(e)}", "error"
    
    def _extract_word(self, file_path: Path) -> tuple[str, str]:
        """Extract text from Word documents."""
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = ["[Table]"]
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_data))
                text_parts.extend(table_text)
            
            return "\n".join(text_parts), "python-docx"
            
        except ImportError:
            return "", "python-docx_not_available"
        except Exception as e:
            return f"Word extraction error: {str(e)}", "error"
    
    def _extract_excel(self, file_path: Path) -> tuple[str, str]:
        """Extract text from Excel files."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(str(file_path), data_only=True)
            sheets_text = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_content = [f"=== Sheet: {sheet_name} ==="]
                
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):
                        sheet_content.append(" | ".join(row_data))
                
                sheets_text.append("\n".join(sheet_content))
            
            return "\n\n".join(sheets_text), "openpyxl"
            
        except ImportError:
            return "", "openpyxl_not_available"
        except Exception as e:
            return f"Excel extraction error: {str(e)}", "error"
    
    def _extract_text(self, file_path: Path) -> tuple[str, str]:
        """Extract text from plain text files."""
        try:
            encodings = ["utf-8", "latin-1", "cp1252"]
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read(), f"text_{encoding}"
                except UnicodeDecodeError:
                    continue
            
            # Fallback with error replacement
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read(), "text_utf8_replace"
                
        except Exception as e:
            return f"Text extraction error: {str(e)}", "error"
    
    def save_outputs(self, result: Dict[str, Any]) -> None:
        """Save result as both .txt and .json files."""
        base_name = Path(result["file_name"]).stem
        doc_id = str(uuid.uuid4())[:8]
        
        # Save .txt file
        txt_file = self.txt_dir / f"{base_name}_{doc_id}.txt"
        with open(txt_file, 'w', encoding='utf-8') as f:
            f.write(f"File: {result['file_path']}\n")
            f.write(f"Processed: {result['processed_at']}\n")
            f.write(f"Method: {result['extraction_method']}\n")
            f.write("=" * 50 + "\n\n")
            f.write(result['extracted_text'])
        
        # Save .json file
        json_file = self.json_dir / f"{base_name}_{doc_id}.json"
        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump(result, f, indent=2, ensure_ascii=False)
        
        print(f"✅ Saved: {result['file_name']} -> {txt_file.name}, {json_file.name}")
    
    def process_directory(self, input_dir: str) -> None:
        """Process all files in a directory."""
        input_path = Path(input_dir)
        
        if not input_path.exists():
            print(f"❌ Input directory not found: {input_dir}")
            return
        
        # Find ALL files with detailed logging
        files_to_process = []
        skipped_files = []
        
        print(f"🔍 Scanning directory: {input_path}")
        
        for root, dirs, files in os.walk(input_path):
            print(f"📂 Checking folder: {root}")
            print(f"   Found {len(files)} files: {files}")
            
            for file in files:
                file_path = Path(root) / file
                
                # Only skip obvious system files
                if file.startswith('.DS_Store') or file.startswith('Thumbs.db'):
                    skipped_files.append(f"System file: {file}")
                    continue
                
                # Skip extremely large files (>500MB)
                try:
                    size_mb = file_path.stat().st_size / (1024 * 1024)
                    if size_mb > 500:
                        skipped_files.append(f"Too large ({size_mb:.1f}MB): {file}")
                        continue
                except:
                    pass
                
                files_to_process.append(file_path)
                print(f"   ✅ Will process: {file}")
        
        print(f"\n📁 Total files found: {len(files_to_process)}")
        print(f"⚠️  Skipped files: {len(skipped_files)}")
        
        if skipped_files:
            for skip in skipped_files[:10]:  # Show first 10 skipped
                print(f"   - {skip}")
        
        if not files_to_process:
            print("❌ No files to process!")
            return
        
        # Process each file
        for file_path in files_to_process:
            print(f"🔄 Processing: {file_path.name}")
            
            result = self.process_file(str(file_path))
            self.save_outputs(result)
            
            # Update stats
            if result["success"]:
                self.stats["processed"] += 1
            else:
                self.stats["failed"] += 1
                print(f"❌ Failed: {file_path.name} - {result['errors']}")
            
            ext = result["file_extension"]
            self.stats["by_type"][ext] = self.stats["by_type"].get(ext, 0) + 1
        
        # Save summary
        summary = {
            "processing_summary": self.stats,
            "processed_at": datetime.now().isoformat(),
            "output_directory": str(self.output_dir.absolute())
        }
        
        with open(self.output_dir / "processing_summary.json", 'w') as f:
            json.dump(summary, f, indent=2)
        
        print(f"\n📊 Summary:")
        print(f"   Processed: {self.stats['processed']}")
        print(f"   Failed: {self.stats['failed']}")
        print(f"   By type: {self.stats['by_type']}")
        print(f"   Output: {self.output_dir}")

def main():
    """Main function."""
    import sys
    
    print("🚀 Simple Document Processor")
    print("=" * 40)
    
    # Check dependencies
    print("🔍 Checking dependencies...")
    if not setup_dependencies():
        print("\n⚠️  Install missing dependencies and try again")
        return
    
    # Get paths
    if len(sys.argv) < 3:
        print("\nUsage: python simple_processor.py INPUT_DIR OUTPUT_DIR")
        print("\nExample:")
        print("python simple_processor.py /path/to/documents /path/to/output")
        return
    
    input_dir = sys.argv[1]
    output_dir = sys.argv[2]
    
    print(f"📁 Input: {input_dir}")
    print(f"📁 Output: {output_dir}")
    
    # Process documents
    processor = SimpleDocumentProcessor(output_dir)
    processor.process_directory(input_dir)
    
    print(f"\n✅ Processing complete!")

if __name__ == "__main__":
    main()