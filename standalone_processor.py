#!/usr/bin/env python3
"""
Standalone document processor script with fixed imports.
"""
import os
import sys
import logging
import zipfile
import tempfile
import shutil
import uuid
import json
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple
from dataclasses import dataclass
from concurrent.futures import ThreadPoolExecutor, as_completed

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(f"dataset_processing_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log")
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class Document:
    """Represents a processed document in the system."""
    doc_id: str
    file_name: str
    file_extension: str
    content_type: List[str] = None
    created_at: Optional[datetime] = None
    author: Optional[str] = None
    text_content: str = ""
    metadata: Dict[str, Any] = None
    persuasion_tags: List[str] = None
    box_folder: Optional[str] = None
    chunks: List[Dict[str, Any]] = None
    
    def __post_init__(self):
        if self.content_type is None:
            self.content_type = []
        if self.metadata is None:
            self.metadata = {}
        if self.persuasion_tags is None:
            self.persuasion_tags = []
        if self.chunks is None:
            self.chunks = []

@dataclass
class DocumentChunk:
    """Represents a chunk of text from a document."""
    chunk_id: str
    doc_id: str
    text_chunk: str
    embedding: Optional[List[float]] = None
    tag_context: List[str] = None
    
    def __post_init__(self):
        if self.tag_context is None:
            self.tag_context = []

@dataclass
class DocumentPair:
    """Represents a pair of metadata and content documents."""
    content_doc: Optional[Document] = None
    metadata_doc: Optional[Document] = None
    zip_file: str = ""
    category: str = ""
    base_name: str = ""
    extraction_path: str = ""
    errors: List[str] = None
    
    def __post_init__(self):
        if self.errors is None:
            self.errors = []

class DocumentClassifier:
    """Document classification and tagging system."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize classifier with configuration."""
        self.config = config or {}
        self.enabled = self.config.get("enabled", True)
    
    def classify(self, document: Document) -> Document:
        """Classify document and add metadata tags."""
        
        if not self.enabled:
            return document
        
        # Content-based classification
        document = self._classify_content_type(document)
        
        # Add persuasion tags
        document = self._add_persuasion_tags(document)
        
        # Add quality metrics
        document = self._assess_quality(document)
        
        return document
    
    def _classify_content_type(self, document: Document) -> Document:
        """Classify document content type based on text analysis."""
        
        if not document.text_content:
            return document
        
        text_lower = document.text_content.lower()
        
        # Academic/research content
        if any(term in text_lower for term in ['abstract', 'methodology', 'conclusion', 'references']):
            if 'academic' not in document.content_type:
                document.content_type.append('academic')
        
        # Legal content
        if any(term in text_lower for term in ['whereas', 'jurisdiction', 'plaintiff', 'defendant']):
            if 'legal' not in document.content_type:
                document.content_type.append('legal')
        
        # Financial content
        if any(term in text_lower for term in ['revenue', 'profit', 'financial', 'investment']):
            if 'financial' not in document.content_type:
                document.content_type.append('financial')
        
        # Email content
        if any(term in text_lower for term in ['subject:', 'from:', 'to:', 'dear']):
            if 'email' not in document.content_type:
                document.content_type.append('email')
        
        return document
    
    def _add_persuasion_tags(self, document: Document) -> Document:
        """Add basic persuasion strategy tags."""
        
        if not document.text_content:
            return document
        
        persuasion_tags = []
        text_lower = document.text_content.lower()
        
        # Authority indicators
        if any(term in text_lower for term in ['expert', 'authority', 'certified', 'official']):
            persuasion_tags.append('authority')
        
        # Social proof indicators  
        if any(term in text_lower for term in ['popular', 'majority', 'everyone', 'most people']):
            persuasion_tags.append('social_proof')
        
        # Urgency indicators
        if any(term in text_lower for term in ['urgent', 'deadline', 'limited time', 'act now']):
            persuasion_tags.append('urgency')
        
        document.persuasion_tags.extend(persuasion_tags)
        
        return document
    
    def _assess_quality(self, document: Document) -> Document:
        """Assess document quality and add metrics."""
        
        quality_score = 1.0
        quality_issues = []
        
        # Text content quality
        if not document.text_content.strip():
            quality_score -= 0.5
            quality_issues.append('no_text_content')
        elif len(document.text_content) < 100:
            quality_score -= 0.2
            quality_issues.append('very_short_content')
        
        # Extraction quality
        if 'extraction_error' in document.metadata:
            quality_score -= 0.3
            quality_issues.append('extraction_error')
        
        # Store quality metrics
        document.metadata['quality_score'] = max(0.0, quality_score)
        if quality_issues:
            document.metadata['quality_issues'] = quality_issues
        
        return document

class ContentAwareChunker:
    """Chunks document text based on content type."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize with configuration."""
        self.config = config or {}
        
        # Default chunking parameters
        self.default_chunk_size = self.config.get("default_chunk_size", 500)
        self.default_chunk_overlap = self.config.get("default_chunk_overlap", 50)
        
        # Content-specific parameters
        self.email_chunk_size = self.config.get("email_chunk_size", 300)
        self.email_chunk_overlap = self.config.get("email_chunk_overlap", 30)
        
        self.presentation_chunk_size = self.config.get("presentation_chunk_size", 250)
        self.presentation_chunk_overlap = self.config.get("presentation_chunk_overlap", 25)
        
        self.image_chunk_size = self.config.get("image_chunk_size", 400)
        self.image_chunk_overlap = self.config.get("image_chunk_overlap", 40)
    
    def chunk_document(self, document: Document) -> Document:
        """Chunk document based on its content type."""
        # Determine chunking strategy based on content type
        if "email" in document.content_type:
            document = self._chunk_email(document)
        elif "presentation" in document.content_type:
            document = self._chunk_presentation(document)
        elif "image" in document.content_type:
            document = self._chunk_image(document)
        else:
            document = self._chunk_regular_document(document)
            
        return document
    
    def _chunk_regular_document(self, document: Document) -> Document:
        """Chunk a regular document with paragraph awareness."""
        import re
        
        text = document.text_content
        chunks = []
        
        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        
        current_chunk = ""
        
        for paragraph in paragraphs:
            # If adding this paragraph exceeds chunk size, create a new chunk
            if len(current_chunk) + len(paragraph) > self.default_chunk_size:
                if current_chunk:
                    chunk = DocumentChunk(
                        chunk_id=str(uuid.uuid4()),
                        doc_id=document.doc_id,
                        text_chunk=current_chunk.strip(),
                        tag_context=document.persuasion_tags
                    )
                    chunks.append(chunk.__dict__)
                
                # Start new chunk with overlap from end of previous chunk
                if current_chunk and self.default_chunk_overlap > 0:
                    # Add overlap from previous chunk
                    words = current_chunk.split()
                    overlap_word_count = min(len(words), self.default_chunk_overlap)
                    overlap_text = ' '.join(words[-overlap_word_count:])
                    current_chunk = overlap_text + ' ' + paragraph
                else:
                    current_chunk = paragraph
            else:
                # Add paragraph to current chunk
                if current_chunk:
                    current_chunk += '\n\n' + paragraph
                else:
                    current_chunk = paragraph
        
        # Add final chunk if not empty
        if current_chunk:
            chunk = DocumentChunk(
                chunk_id=str(uuid.uuid4()),
                doc_id=document.doc_id,
                text_chunk=current_chunk.strip(),
                tag_context=document.persuasion_tags
            )
            chunks.append(chunk.__dict__)
            
        # Update document with chunks
        document.chunks = chunks
        return document
    
    def _chunk_email(self, document: Document) -> Document:
        """Chunk an email with header and thread awareness."""
        import re
        
        # Get full text
        text = document.text_content
        chunks = []
        
        # Split into email thread parts if available
        # Use thread markers to detect parts
        thread_markers = [
            r'----+ ?Original Message ?----+',
            r'----+ ?Forwarded Message ?----+',
            r'On .* wrote:',
            r'From:.*Sent:.*To:.*Subject:'
        ]
        
        # Combine markers into single regex
        combined_regex = '|'.join(thread_markers)
        
        # Split by thread markers
        thread_parts = re.split(combined_regex, text)
        
        # Create chunk for each thread part
        for i, part in enumerate(thread_parts):
            # Skip empty parts
            if not part.strip():
                continue
                
            # Create chunk
            chunk = DocumentChunk(
                chunk_id=str(uuid.uuid4()),
                doc_id=document.doc_id,
                text_chunk=part.strip(),
                tag_context=document.persuasion_tags + [f"email_part_{i}"]
            )
            chunks.append(chunk.__dict__)
            
        # If we couldn't split into thread parts, chunk by size
        if not chunks:
            # Fall back to regular chunking
            return self._chunk_regular_document(document)
            
        # Update document with chunks
        document.chunks = chunks
        return document
    
    def _chunk_presentation(self, document: Document) -> Document:
        """Chunk a presentation with slide awareness."""
        import re
        
        text = document.text_content
        chunks = []
        
        # Split by slide markers
        slide_markers = re.finditer(r'---\s*SLIDE\s+(\d+)\s*---', text)
        slide_positions = [(m.start(), m.group(1)) for m in slide_markers]
        
        if not slide_positions:
            # Fall back to regular chunking if no slide markers
            return self._chunk_regular_document(document)
        
        # Process each slide as a chunk
        for i in range(len(slide_positions)):
            start_pos = slide_positions[i][0]
            slide_num = slide_positions[i][1]
            
            # Determine end position
            if i < len(slide_positions) - 1:
                end_pos = slide_positions[i+1][0]
            else:
                end_pos = len(text)
                
            # Extract slide text
            slide_text = text[start_pos:end_pos].strip()
            
            # Create chunk for slide
            chunk = DocumentChunk(
                chunk_id=str(uuid.uuid4()),
                doc_id=document.doc_id,
                text_chunk=slide_text,
                tag_context=document.persuasion_tags + [f"slide_{slide_num}"]
            )
            chunks.append(chunk.__dict__)
        
        # Update document with chunks
        document.chunks = chunks
        return document
    
    def _chunk_image(self, document: Document) -> Document:
        """Chunk an image document with content awareness."""
        import re
        
        text = document.text_content
        chunks = []
        
        if not text.strip():
            # Empty text - create empty chunk
            chunk = DocumentChunk(
                chunk_id=str(uuid.uuid4()),
                doc_id=document.doc_id,
                text_chunk="",
                tag_context=document.persuasion_tags + ["image"]
            )
            chunks.append(chunk.__dict__)
            document.chunks = chunks
            return document
        
        # For images, we might have OCR text from different visual elements
        # Split by visual element markers if present
        if "[Visual Element" in text:
            elements = re.split(r'\[Visual Element[^\]]*\]:', text)
            
            # First element might be general text
            if elements[0].strip():
                chunk = DocumentChunk(
                    chunk_id=str(uuid.uuid4()),
                    doc_id=document.doc_id,
                    text_chunk=elements[0].strip(),
                    tag_context=document.persuasion_tags + ["image", "general"]
                )
                chunks.append(chunk.__dict__)
            
            # Process each visual element
            for i, element in enumerate(elements[1:], 1):
                if element.strip():
                    chunk = DocumentChunk(
                        chunk_id=str(uuid.uuid4()),
                        doc_id=document.doc_id,
                        text_chunk=element.strip(),
                        tag_context=document.persuasion_tags + ["image", f"visual_element_{i}"]
                    )
                    chunks.append(chunk.__dict__)
        else:
            # No visual elements marked - create single chunk
            chunk = DocumentChunk(
                chunk_id=str(uuid.uuid4()),
                doc_id=document.doc_id,
                text_chunk=text.strip(),
                tag_context=document.persuasion_tags + ["image"]
            )
            chunks.append(chunk.__dict__)
        
        document.chunks = chunks
        return document

class ZipDatasetProcessor:
    """Processor for zip-based datasets with paired files."""
    
    def __init__(self, config: Dict[str, Any]):
        """Initialize the zip dataset processor."""
        self.config = config
        
        # Initialize core components
        self.chunker = ContentAwareChunker(config.get("chunker", {}))
        self.classifier = DocumentClassifier(config.get("classifier", {}))
        
        # Processing statistics
        self.stats = {
            "zip_files_processed": 0,
            "zip_files_failed": 0,
            "document_pairs_created": 0,
            "content_docs_processed": 0,
            "metadata_docs_processed": 0,
            "by_category": {},
            "errors": []
        }
        
        # Temporary extraction directory
        self.temp_dir = tempfile.mkdtemp(prefix="zip_dataset_")
        logger.info(f"Using temporary directory: {self.temp_dir}")
    
    def process_dataset(self, input_root: str, output_dir: Optional[str] = None) -> List[DocumentPair]:
        """Process the entire dataset structure."""
        
        logger.info(f"Starting dataset processing from: {input_root}")
        
        if not os.path.exists(input_root):
            raise ValueError(f"Input directory does not exist: {input_root}")
        
        # Discover all zip files organized by category
        zip_files_by_category = self._discover_zip_files(input_root)
        
        total_zips = sum(len(files) for files in zip_files_by_category.values())
        logger.info(f"Found {total_zips} zip files across {len(zip_files_by_category)} categories")
        
        # Process all zip files
        all_document_pairs = []
        
        for category, zip_files in zip_files_by_category.items():
            logger.info(f"Processing {len(zip_files)} zip files in category: {category}")
            
            category_pairs = self._process_category(category, zip_files)
            all_document_pairs.extend(category_pairs)
            
            # Update statistics
            self.stats["by_category"][category] = {
                "zip_files": len(zip_files),
                "document_pairs": len(category_pairs),
                "successful_pairs": len([p for p in category_pairs if p.content_doc and p.metadata_doc])
            }
        
        # Save processing results if output directory specified
        if output_dir:
            self._save_results(all_document_pairs, output_dir)
        
        logger.info(f"Dataset processing complete. Processed {len(all_document_pairs)} document pairs")
        
        return all_document_pairs
    
    def _discover_zip_files(self, input_root: str) -> Dict[str, List[str]]:
        """Discover all zip files organized by category."""
        
        zip_files_by_category = {}
        
        # Walk through the directory structure
        for category_dir in os.listdir(input_root):
            category_path = os.path.join(input_root, category_dir)
            
            if not os.path.isdir(category_path):
                continue
            
            # Skip hidden directories and files
            if category_dir.startswith('.'):
                continue
                
            zip_files = []
            
            # Find all zip files in this category
            for file_name in os.listdir(category_path):
                if file_name.endswith('.zip'):
                    zip_path = os.path.join(category_path, file_name)
                    zip_files.append(zip_path)
            
            if zip_files:
                zip_files_by_category[category_dir] = sorted(zip_files)
                logger.info(f"Category '{category_dir}': {len(zip_files)} zip files")
        
        return zip_files_by_category
    
    def _process_category(self, category: str, zip_files: List[str]) -> List[DocumentPair]:
        """Process all zip files in a category."""
        
        document_pairs = []
        
        # Process zip files in parallel
        max_workers = self.config.get("max_workers", 4)
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all zip processing jobs
            future_to_zip = {
                executor.submit(self._process_zip_file, zip_file, category): zip_file
                for zip_file in zip_files
            }
            
            # Collect results
            for future in as_completed(future_to_zip):
                zip_file = future_to_zip[future]
                
                try:
                    document_pair = future.result()
                    document_pairs.append(document_pair)
                    
                    if document_pair.content_doc and document_pair.metadata_doc:
                        logger.info(f"Successfully processed zip: {os.path.basename(zip_file)}")
                        self.stats["zip_files_processed"] += 1
                    else:
                        logger.warning(f"Partial processing for zip: {os.path.basename(zip_file)}")
                        if document_pair.errors:
                            logger.warning(f"Errors: {'; '.join(document_pair.errors)}")
                        
                except Exception as e:
                    logger.error(f"Failed to process zip {zip_file}: {e}")
                    self.stats["zip_files_failed"] += 1
                    self.stats["errors"].append(f"{zip_file}: {str(e)}")
                    
                    # Create error document pair
                    document_pairs.append(DocumentPair(
                        zip_file=zip_file,
                        category=category,
                        base_name=os.path.splitext(os.path.basename(zip_file))[0],
                        errors=[str(e)]
                    ))
        
        return document_pairs
    
    def _process_zip_file(self, zip_file: str, category: str) -> DocumentPair:
        """Process a single zip file containing file pair."""
        
        base_name = os.path.splitext(os.path.basename(zip_file))[0]
        
        # Create document pair
        document_pair = DocumentPair(
            zip_file=zip_file,
            category=category,
            base_name=base_name
        )
        
        try:
            # Extract zip file
            extraction_path = self._extract_zip_file(zip_file, base_name)
            document_pair.extraction_path = extraction_path
            
            # Find all processable files
            processable_files = self._find_content_files(extraction_path)
            
            if len(processable_files) == 0:
                document_pair.errors.append(f"No processable files found in ZIP")
                return document_pair
            
            # Identify metadata and content files based on category
            if category == "Presentations":
                metadata_file, content_file = self._identify_presentation_pair(processable_files, base_name)
            elif category == "Images":
                metadata_file, content_file = self._identify_image_pair(processable_files, base_name)
            elif category == "Spreadsheets":
                metadata_file, content_file = self._identify_spreadsheet_pair(processable_files, base_name)
            else:
                # Default to PDF pair identification
                pdf_files = [f for f in processable_files if f.lower().endswith('.pdf')]
                if len(pdf_files) >= 2:
                    metadata_file, content_file = self._identify_pdf_pair(pdf_files, base_name)
                else:
                    # Fallback for mixed file types
                    metadata_file, content_file = self._identify_mixed_pair(processable_files, base_name)
            
            # Process metadata file
            if metadata_file:
                metadata_result = self._process_file(metadata_file, category, "metadata", base_name)
                if metadata_result:
                    document_pair.metadata_doc = metadata_result
                    self.stats["metadata_docs_processed"] += 1
                else:
                    document_pair.errors.append(f"Failed to process metadata file: {os.path.basename(metadata_file)}")
            else:
                document_pair.errors.append("No metadata file found")
            
            # Process content file
            if content_file:
                content_result = self._process_file(content_file, category, "content", base_name)
                if content_result:
                    document_pair.content_doc = content_result
                    self.stats["content_docs_processed"] += 1
                else:
                    document_pair.errors.append(f"Failed to process content file: {os.path.basename(content_file)}")
            else:
                document_pair.errors.append("No content file found")
            
            # Link the documents if both successful
            if document_pair.content_doc and document_pair.metadata_doc:
                self._link_document_pair(document_pair)
                self.stats["document_pairs_created"] += 1
            
        except Exception as e:
            document_pair.errors.append(str(e))
            logger.error(f"Error processing zip {zip_file}: {e}")
            import traceback
            traceback.print_exc()
        
        finally:
            # Clean up extracted files
            if document_pair.extraction_path and os.path.exists(document_pair.extraction_path):
                try:
                    shutil.rmtree(document_pair.extraction_path)
                except Exception as e:
                    logger.warning(f"Failed to clean up {document_pair.extraction_path}: {e}")
        
        return document_pair
    
    def _extract_zip_file(self, zip_file: str, base_name: str) -> str:
        """Extract zip file to temporary directory."""
        
        extraction_path = os.path.join(self.temp_dir, f"{base_name}_{uuid.uuid4().hex[:8]}")
        os.makedirs(extraction_path, exist_ok=True)
        
        with zipfile.ZipFile(zip_file, 'r') as zip_ref:
            zip_ref.extractall(extraction_path)
        
        # Handle nested ZIPs
        self._handle_nested_zips(extraction_path)
        
        return extraction_path
    
    def _handle_nested_zips(self, directory: str):
        """Extract any nested ZIP files."""
        for root, _, files in os.walk(directory):
            for file in files:
                if file.endswith('.zip'):
                    nested_zip = os.path.join(root, file)
                    extract_dir = os.path.join(root, f"{os.path.splitext(file)[0]}_extracted")
                    
                    try:
                        with zipfile.ZipFile(nested_zip, 'r') as zip_ref:
                            zip_ref.extractall(extract_dir)
                        os.remove(nested_zip)
                        logger.info(f"Extracted nested ZIP: {file}")
                    except Exception as e:
                        logger.warning(f"Failed to extract nested ZIP {file}: {e}")
    
    def _find_content_files(self, extraction_path: str) -> List[str]:
        """Find all processable content files."""
        processable_files = []
        supported_extensions = ['.pdf', '.xlsx', '.xls', '.csv', '.docx', '.doc', '.pptx', '.ppt', 
                              '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.txt', '.md']
        
        for root, _, files in os.walk(extraction_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_ext = os.path.splitext(file_path)[1].lower()
                if file_ext in supported_extensions:
                    processable_files.append(file_path)
        
        return processable_files
    
    def _identify_pdf_pair(self, pdf_files: List[str], base_name: str) -> Tuple[Optional[str], Optional[str]]:
        """Identify which PDF is metadata and which is content."""
        
        metadata_pdf = None
        content_pdf = None
        
        for pdf_file in pdf_files:
            file_name = os.path.basename(pdf_file)
            name_without_ext = os.path.splitext(file_name)[0]
            
            # Check if this is the metadata file (contains -info suffix)
            if name_without_ext.endswith('-info') or 'info' in name_without_ext.lower():
                metadata_pdf = pdf_file
            else:
                # Check if this matches the base name or is the main content
                if name_without_ext == base_name or name_without_ext.startswith(base_name):
                    content_pdf = pdf_file
        
        # If we couldn't identify by naming convention, use file size heuristic
        if not metadata_pdf or not content_pdf:
            # Sort by file size - metadata is usually smaller
            pdf_files_with_size = [(f, os.path.getsize(f)) for f in pdf_files]
            pdf_files_with_size.sort(key=lambda x: x[1])
            
            # Smaller file is likely metadata
            metadata_pdf = pdf_files_with_size[0][0]
            # Larger file is likely content
            content_pdf = pdf_files_with_size[-1][0]
        
        return metadata_pdf, content_pdf
    
    def _identify_presentation_pair(self, files: List[str], base_name: str) -> Tuple[Optional[str], Optional[str]]:
        """Identify content and metadata files for presentations."""
        
        # Look for presentation files (content) and PDF info files (metadata)
        presentation_files = [f for f in files if f.lower().endswith(('.pptx', '.ppt'))]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        
        content_file = None
        metadata_file = None
        
        if presentation_files:
            content_file = presentation_files[0]  # Use first presentation file as content
        
        # Find metadata file (typically PDF with 'info' in name)
        for pdf_file in pdf_files:
            if 'info' in os.path.basename(pdf_file).lower():
                metadata_file = pdf_file
                break
        
        # If no explicit metadata file found, use file size heuristic
        if not metadata_file and pdf_files:
            # Sort PDFs by size - metadata typically smaller
            pdf_files_with_size = [(f, os.path.getsize(f)) for f in pdf_files]
            pdf_files_with_size.sort(key=lambda x: x[1])
            metadata_file = pdf_files_with_size[0][0]
        
        return metadata_file, content_file
    
    def _identify_image_pair(self, files: List[str], base_name: str) -> Tuple[Optional[str], Optional[str]]:
        """Identify content and metadata files for images."""
        
        # Look for image files (content) and PDF info files (metadata)
        image_files = [f for f in files if any(f.lower().endswith(ext) for ext in 
                        ['.png', '.jpg', '.jpeg', '.tiff', '.bmp'])]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        
        content_file = None
        metadata_file = None
        
        if image_files:
            content_file = image_files[0]  # Use first image file as content
        
        # Find metadata file (typically PDF with 'info' in name)
        for pdf_file in pdf_files:
            if 'info' in os.path.basename(pdf_file).lower():
                metadata_file = pdf_file
                break
        
        # If no explicit metadata file found, use file size heuristic
        if not metadata_file and pdf_files:
            metadata_file = pdf_files[0]
        
        return metadata_file, content_file
    
    def _identify_spreadsheet_pair(self, files: List[str], base_name: str) -> Tuple[Optional[str], Optional[str]]:
        """Identify content and metadata for spreadsheet files."""
        # Separate by file type
        excel_files = [f for f in files if any(f.lower().endswith(ext) for ext in ['.xlsx', '.xls'])]
        csv_files = [f for f in files if f.lower().endswith('.csv')]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        
        content_file = None
        metadata_file = None
        
        # Content is usually the Excel/CSV file
        if excel_files:
            content_file = excel_files[0]
        elif csv_files:
            content_file = csv_files[0]
        
        # Metadata is usually the PDF info file
        for pdf in pdf_files:
            if 'info' in os.path.basename(pdf).lower():
                metadata_file = pdf
                break
        
        # If no metadata PDF, use the smaller file
        if not metadata_file and len(files) >= 2:
            files_by_size = sorted(files, key=lambda f: os.path.getsize(f))
            metadata_file = files_by_size[0]
            if not content_file:
                content_file = files_by_size[-1]
        
        return metadata_file, content_file
    
    def _identify_mixed_pair(self, files: List[str], base_name: str) -> Tuple[Optional[str], Optional[str]]:
        """Identify pairs from mixed file types."""
        if len(files) == 1:
            # Single file - treat as content
            return None, files[0]
        
        # Sort by size - metadata usually smaller
        files_by_size = sorted(files, key=lambda f: os.path.getsize(f))
        
        # Check for 'info' pattern
        metadata_file = None
        for f in files:
            if 'info' in os.path.basename(f).lower():
                metadata_file = f
                break
        
        if metadata_file:
            content_file = next((f for f in files if f != metadata_file), None)
        else:
            # Use size heuristic
            metadata_file = files_by_size[0]
            content_file = files_by_size[-1]
        
        return metadata_file, content_file
    
    def _process_file(self, file_path: str, category: str, doc_type: str, base_name: str) -> Optional[Document]:
        """Process a single file."""
        try:
            # Extract text based on file type
            file_extension = os.path.splitext(file_path)[1][1:].lower()
            doc_id = str(uuid.uuid4())
            file_name = os.path.basename(file_path)
            
            # Create base document
            document = Document(
                doc_id=doc_id,
                file_name=file_name,
                file_extension=file_extension,
                content_type=[file_extension],
                created_at=datetime.now(),
                metadata={
                    "dataset_category": category,
                    "document_type": doc_type,
                    "base_name": base_name,
                    "source_zip": f"{base_name}.zip"
                }
            )
            
            # Extract text based on file type
            if file_extension == "pdf":
                document = self._extract_pdf_text(file_path, document)
            elif file_extension in ["pptx", "ppt"]:
                document = self._extract_presentation_text(file_path, document)
            elif file_extension in ["xlsx", "xls", "csv"]:
                document = self._extract_spreadsheet_text(file_path, document)
            elif file_extension in ["docx", "doc"]:
                document = self._extract_word_text(file_path, document)
            elif file_extension in ["png", "jpg", "jpeg", "tiff", "bmp"]:
                document = self._extract_image_text(file_path, document)
            else:
                # Default to text extraction
                document = self._extract_text_file(file_path, document)
            
            # Update content type
            if doc_type == "metadata":
                if "metadata" not in document.content_type:
                    document.content_type.append("metadata")
            
            # Add content type based on category
            if category == "Presentations":
                if "presentation" not in document.content_type:
                    document.content_type.append("presentation")
            elif category == "Images":
                if "image" not in document.content_type:
                    document.content_type.append("image")
            elif category == "Spreadsheets":
                if "spreadsheet" not in document.content_type:
                    document.content_type.append("spreadsheet")
            elif category == "Emails":
                if "email" not in document.content_type:
                    document.content_type.append("email")
            
            # Classify document
            document = self.classifier.classify(document)
            
            # Chunk document
            document = self.chunker.chunk_document(document)
            
            return document
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            import traceback
            traceback.print_exc()
            return None
    
    def _extract_pdf_text(self, file_path: str, document: Document) -> Document:
        """Extract text from PDF file."""
        try:
            import fitz  # PyMuPDF
            
            with fitz.open(file_path) as pdf:
                text = ""
                for page in pdf:
                    text += page.get_text() + "\n\n"
                
                # Add metadata
                document.metadata.update({
                    "page_count": len(pdf),
                    "title": pdf.metadata.get("title", ""),
                    "author": pdf.metadata.get("author", "")
                })
                
                # Set author if available
                if pdf.metadata.get("author"):
                    document.author = pdf.metadata.get("author")
                
                # Set created_at if available
                if pdf.metadata.get("creationDate"):
                    # PyMuPDF dates need special parsing - this is simplified
                    try:
                        # Convert PDF date format to datetime
                        date_str = pdf.metadata.get("creationDate")
                        # Remove D: prefix and Z suffix if present
                        if date_str.startswith("D:"):
                            date_str = date_str[2:]
                        if date_str.endswith("Z"):
                            date_str = date_str[:-1]
                        # Format: YYYYMMDDHHmmSS
                        document.created_at = datetime.strptime(date_str[:14], "%Y%m%d%H%M%S")
                    except:
                        pass
            
            document.text_content = text.strip()
            return document
            
        except ImportError:
            logger.warning("PyMuPDF not available - using fallback text extraction")
            
            # Fallback to simple text extraction
            try:
                with open(file_path, 'rb') as f:
                    document.text_content = f"[PDF Content - Unable to extract text. Install PyMuPDF for better results.]"
                    document.metadata["extraction_error"] = "PyMuPDF not available"
                return document
            except Exception as e:
                document.metadata["extraction_error"] = str(e)
                return document
                
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _extract_presentation_text(self, file_path: str, document: Document) -> Document:
        """Extract text from presentation file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"--- SLIDE {i+1} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                slides_text.append("\n".join(slide_text))
            
            document.text_content = "\n\n".join(slides_text)
            document.metadata["slide_count"] = len(prs.slides)
            document.content_type.append("presentation")
            return document
            
        except ImportError:
            logger.warning("python-pptx not available - using fallback extraction")
            document.text_content = f"[Presentation Content - Unable to extract text. Install python-pptx for better results.]"
            document.metadata["extraction_error"] = "python-pptx not available"
            return document
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _extract_spreadsheet_text(self, file_path: str, document: Document) -> Document:
        """Extract text from spreadsheet file."""
        try:
            if file_path.endswith('.csv'):
                import csv
                
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    csv_reader = csv.reader(f)
                    rows = list(csv_reader)
                    
                    text_content = []
                    for row in rows[:100]:  # Limit to first 100 rows
                        text_content.append("\t".join(row))
                    
                    document.text_content = "\n".join(text_content)
                    document.metadata["row_count"] = len(rows)
                    document.metadata["column_count"] = len(rows[0]) if rows else 0
                    document.content_type.append("spreadsheet")
                    return document
            else:
                import openpyxl
                
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                text_content = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    text_content.append(f"=== SHEET: {sheet_name} ===")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                            else:
                                row_text.append("")
                        text_content.append("\t".join(row_text))
                
                document.text_content = "\n".join(text_content)
                document.metadata["sheet_count"] = len(workbook.sheetnames)
                document.content_type.append("spreadsheet")
                return document
                
        except ImportError:
            logger.warning("Excel libraries not available - using fallback extraction")
            document.text_content = f"[Spreadsheet Content - Unable to extract text. Install openpyxl for better results.]"
            document.metadata["extraction_error"] = "Excel libraries not available"
            return document
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _extract_word_text(self, file_path: str, document: Document) -> Document:
        """Extract text from Word document."""
        try:
            from docx import Document as DocxDocument
            
            doc = DocxDocument(file_path)
            paragraphs = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            document.text_content = "\n\n".join(paragraphs)
            document.metadata["paragraph_count"] = len(paragraphs)
            document.content_type.append("document")
            return document
            
        except ImportError:
            logger.warning("python-docx not available - using fallback extraction")
            document.text_content = f"[Word Document Content - Unable to extract text. Install python-docx for better results.]"
            document.metadata["extraction_error"] = "python-docx not available"
            return document
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _extract_image_text(self, file_path: str, document: Document) -> Document:
        """Extract text from image file using OCR."""
        try:
            from PIL import Image
            import pytesseract
            
            with Image.open(file_path) as img:
                # Basic image metadata
                document.metadata.update({
                    "width": img.width,
                    "height": img.height,
                    "format": img.format,
                    "mode": img.mode
                })
                
                # Perform OCR
                text = pytesseract.image_to_string(img)
                document.text_content = text.strip()
                document.content_type.append("image")
                
                # Simple image type detection
                if "chart" in document.text_content.lower() or "graph" in document.text_content.lower():
                    document.metadata["image_type"] = "chart"
                elif "diagram" in document.text_content.lower() or "flow" in document.text_content.lower():
                    document.metadata["image_type"] = "diagram"
                else:
                    document.metadata["image_type"] = "general_image"
                
                return document
                
        except ImportError:
            logger.warning("PIL/Tesseract not available - using fallback extraction")
            document.text_content = f"[Image Content - Unable to extract text. Install PIL and pytesseract for OCR.]"
            document.metadata["extraction_error"] = "OCR libraries not available"
            return document
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _extract_text_file(self, file_path: str, document: Document) -> Document:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                document.text_content = f.read()
            return document
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
            return document
    
    def _link_document_pair(self, document_pair: DocumentPair):
        """Create bidirectional links between paired documents."""
        
        if not document_pair.content_doc or not document_pair.metadata_doc:
            return
        
        # Add relationship metadata
        document_pair.content_doc.metadata.update({
            "metadata_document_id": document_pair.metadata_doc.doc_id,
            "has_metadata_document": True
        })
        
        document_pair.metadata_doc.metadata.update({
            "content_document_id": document_pair.content_doc.doc_id,
            "describes_document": True
        })
        
        # Add to content types
        if "paired" not in document_pair.content_doc.content_type:
            document_pair.content_doc.content_type.append("paired")
        
        if "metadata" not in document_pair.metadata_doc.content_type:
            document_pair.metadata_doc.content_type.append("metadata")
    
    def _save_results(self, document_pairs: List[DocumentPair], output_dir: str):
        """Save processing results to output directory."""
        
        os.makedirs(output_dir, exist_ok=True)
        
        # Save summary statistics
        summary = {
            "processing_timestamp": datetime.now().isoformat(),
            "total_document_pairs": len(document_pairs),
            "successful_pairs": len([p for p in document_pairs if p.content_doc and p.metadata_doc]),
            "statistics": self.stats,
            "document_pairs": []
        }
        
        # Save detailed results
        for pair in document_pairs:
            pair_data = {
                "base_name": pair.base_name,
                "category": pair.category,
                "zip_file": pair.zip_file,
                "errors": pair.errors
            }
            
            # Add content document info
            if pair.content_doc:
                pair_data["content_document"] = {
                    "doc_id": pair.content_doc.doc_id,
                    "file_name": pair.content_doc.file_name,
                    "content_types": pair.content_doc.content_type,
                    "chunks_count": len(pair.content_doc.chunks)
                }
                
                # Save content text to separate file
                content_file = os.path.join(output_dir, f"{pair.base_name}_content.txt")
                with open(content_file, 'w', encoding='utf-8') as f:
                    f.write(f"Content Type: {', '.join(pair.content_doc.content_type)}\n")
                    f.write(f"File: {pair.content_doc.file_name}\n")
                    f.write("=" * 50 + "\n\n")
                    f.write(pair.content_doc.text_content)
            
            # Add metadata document info
            if pair.metadata_doc:
                pair_data["metadata_document"] = {
                    "doc_id": pair.metadata_doc.doc_id,
                    "file_name": pair.metadata_doc.file_name,
                    "content_types": pair.metadata_doc.content_type,
                    "chunks_count": len(pair.metadata_doc.chunks)
                }
                
                # Save metadata text to separate file
                metadata_file = os.path.join(output_dir, f"{pair.base_name}_metadata.txt")
                with open(metadata_file, 'w', encoding='utf-8') as f:
                    f.write(f"Content Type: {', '.join(pair.metadata_doc.content_type)}\n")
                    f.write(f"File: {pair.metadata_doc.file_name}\n")
                    f.write("=" * 50 + "\n\n")
                    f.write(pair.metadata_doc.text_content)
            
            summary["document_pairs"].append(pair_data)
        
        # Save to file
        summary_file = os.path.join(output_dir, f"dataset_processing_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
        with open(summary_file, 'w', encoding='utf-8') as f:
            json.dump(summary, f, indent=2)
        
        logger.info(f"Results saved to: {summary_file}")
    
    def get_statistics(self) -> Dict[str, Any]:
        """Get processing statistics."""
        return self.stats.copy()
    
    def cleanup(self):
        """Clean up temporary files."""
        if os.path.exists(self.temp_dir):
            try:
                shutil.rmtree(self.temp_dir)
                logger.info(f"Cleaned up temporary directory: {self.temp_dir}")
            except Exception as e:
                logger.warning(f"Failed to clean up temp directory: {e}")

def main():
    """Main entry point."""
    
    logger.info("Starting standalone document processor")
    
    # Configuration
    config = {
        "max_workers": 4,
        "chunker": {
            "default_chunk_size": 500,
            "default_chunk_overlap": 50,
            "email_chunk_size": 300,
            "presentation_chunk_size": 250,
            "image_chunk_size": 400
        },
        "classifier": {
            "enabled": True
        }
    }
    
    # Initialize processor
    processor = ZipDatasetProcessor(config)
    
    try:
        # Get input path
        if len(sys.argv) > 1:
            input_path = sys.argv[1]
        else:
            input_path = os.path.join(os.getcwd(), "data", "input")
            # Create if not exists
            os.makedirs(input_path, exist_ok=True)
            logger.info(f"Using default input path: {input_path}")
        
        # Get output path
        if len(sys.argv) > 2:
            output_path = sys.argv[2]
        else:
            output_path = os.path.join(os.getcwd(), "data", "processed")
            os.makedirs(output_path, exist_ok=True)
            logger.info(f"Using default output path: {output_path}")
        
        # Process the dataset
        document_pairs = processor.process_dataset(input_path, output_path)
        
        # Print statistics
        stats = processor.get_statistics()
        print("\nProcessing Statistics:")
        print(f"Total ZIP files processed: {stats['zip_files_processed']}")
        print(f"Document pairs created: {stats['document_pairs_created']}")
        print(f"Content documents: {stats['content_docs_processed']}")
        print(f"Metadata documents: {stats['metadata_docs_processed']}")
        
        print("\nBy Category:")
        for category, cat_stats in stats['by_category'].items():
            success_rate = (cat_stats['successful_pairs'] / cat_stats['zip_files'] 
                          if cat_stats['zip_files'] > 0 else 0)
            print(f"  {category}: {cat_stats['successful_pairs']}/{cat_stats['zip_files']} "
                  f"({success_rate:.1%})")
        
        print(f"\nResults saved to: {output_path}")
        
    except Exception as e:
        logger.error(f"Processing failed: {e}")
        import traceback
        traceback.print_exc()
        
    finally:
        # Clean up
        processor.cleanup()

if __name__ == "__main__":
    main()