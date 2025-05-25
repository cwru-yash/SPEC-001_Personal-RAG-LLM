# # services/processor/src/pipeline/extractors/office_extractor.py
# import os
# import uuid
# from typing import Dict, Any, List, Optional
# from datetime import datetime
# from pathlib import Path

# from ...models.document import Document

# class OfficeExtractor:
#     """Extractor for Office documents (Excel, Word, PowerPoint)."""
    
#     def __init__(self, config: Dict[str, Any] = None):
#         """Initialize with configuration."""
#         self.config = config or {}
#         self.excel_config = self.config.get("excel", {})
#         self.word_config = self.config.get("word", {})
#         self.powerpoint_config = self.config.get("powerpoint", {})
    
#     def extract(self, file_path: str) -> Document:
#         """Extract text and metadata from Office documents."""
#         file_ext = Path(file_path).suffix[1:].lower()
        
#         if file_ext in ['xlsx', 'xls']:
#             return self._extract_excel(file_path)
#         elif file_ext in ['docx', 'doc']:
#             return self._extract_word(file_path)
#         elif file_ext in ['pptx', 'ppt']:
#             return self._extract_powerpoint(file_path)
#         else:
#             raise ValueError(f"Unsupported Office format: {file_ext}")
    
#     def _extract_excel(self, file_path: str) -> Document:
#         """Extract content from Excel files."""
#         doc_id = str(uuid.uuid4())
#         file_name = os.path.basename(file_path)
        
#         try:
#             import openpyxl
            
#             workbook = openpyxl.load_workbook(file_path, data_only=True)
#             text_content = []
#             sheet_summaries = []
            
#             max_rows = self.excel_config.get("max_rows_per_sheet", 1000)
            
#             for sheet_name in workbook.sheetnames:
#                 sheet = workbook[sheet_name]
#                 sheet_text = [f"=== SHEET: {sheet_name} ==="]
#                 row_count = 0
                
#                 # Extract headers
#                 headers = []
#                 for cell in sheet[1]:
#                     if cell.value is not None:
#                         headers.append(str(cell.value))
                
#                 if headers:
#                     sheet_text.append("Headers: " + " | ".join(headers))
                
#                 # Extract data rows
#                 for row_idx, row in enumerate(sheet.iter_rows(min_row=2, values_only=True)):
#                     if row_idx >= max_rows:
#                         sheet_text.append(f"... (truncated at {max_rows} rows)")
#                         break
                    
#                     row_values = []
#                     for cell in row:
#                         if cell is not None:
#                             row_values.append(str(cell))
                    
#                     if any(row_values):
#                         row_count += 1
#                         if row_idx < 10:  # Show first 10 rows
#                             sheet_text.append(" | ".join(row_values))
                
#                 sheet_text.append(f"Total rows: {row_count}")
#                 text_content.extend(sheet_text)
                
#                 sheet_summaries.append({
#                     "name": sheet_name,
#                     "rows": row_count,
#                     "columns": len(headers) if headers else 0,
#                     "headers": headers
#                 })
            
#             workbook.close()
            
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension=file_ext,
#                 content_type=["spreadsheet", "excel"],
#                 created_at=datetime.now(),
#                 text_content="\n\n".join(text_content),
#                 metadata={
#                     "sheet_count": len(workbook.sheetnames),
#                     "sheets": sheet_summaries,
#                     "extraction_method": "openpyxl",
#                     "include_formulas": self.excel_config.get("include_formulas", False)
#                 }
#             )
            
#         except ImportError:
#             raise ImportError("openpyxl not installed. Install with: pip install openpyxl")
#         except Exception as e:
#             # Create error document
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension="xlsx",
#                 content_type=["spreadsheet", "excel", "error"],
#                 created_at=datetime.now(),
#                 text_content=f"Error extracting Excel content: {str(e)}",
#                 metadata={"extraction_error": str(e)}
#             )
    
#     def _extract_word(self, file_path: str) -> Document:
#         """Extract content from Word documents."""
#         doc_id = str(uuid.uuid4())
#         file_name = os.path.basename(file_path)
        
#         try:
#             from docx import Document as DocxDocument
            
#             doc = DocxDocument(file_path)
#             paragraphs = []
            
#             # Extract paragraphs
#             for paragraph in doc.paragraphs:
#                 if paragraph.text.strip():
#                     paragraphs.append(paragraph.text.strip())
            
#             # Extract tables if configured
#             tables_text = []
#             if self.word_config.get("extract_tables", True):
#                 for table_idx, table in enumerate(doc.tables):
#                     table_text = [f"\n=== TABLE {table_idx + 1} ==="]
#                     for row in table.rows:
#                         row_text = []
#                         for cell in row.cells:
#                             row_text.append(cell.text.strip())
#                         table_text.append(" | ".join(row_text))
#                     tables_text.append("\n".join(table_text))
            
#             # Combine all text
#             all_text = "\n\n".join(paragraphs)
#             if tables_text:
#                 all_text += "\n\n" + "\n\n".join(tables_text)
            
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension="docx",
#                 content_type=["document", "word"],
#                 created_at=datetime.now(),
#                 text_content=all_text,
#                 metadata={
#                     "paragraph_count": len(paragraphs),
#                     "table_count": len(doc.tables),
#                     "extraction_method": "python-docx"
#                 }
#             )
            
#         except ImportError:
#             raise ImportError("python-docx not installed. Install with: pip install python-docx")
#         except Exception as e:
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension="docx",
#                 content_type=["document", "word", "error"],
#                 created_at=datetime.now(),
#                 text_content=f"Error extracting Word content: {str(e)}",
#                 metadata={"extraction_error": str(e)}
#             )
    
#     def _extract_powerpoint(self, file_path: str) -> Document:
#         """Extract content from PowerPoint presentations."""
#         doc_id = str(uuid.uuid4())
#         file_name = os.path.basename(file_path)
        
#         try:
#             from pptx import Presentation
            
#             prs = Presentation(file_path)
#             slides_content = []
#             slide_summaries = []
            
#             for slide_idx, slide in enumerate(prs.slides):
#                 slide_text = [f"=== SLIDE {slide_idx + 1} ==="]
#                 slide_shapes = []
                
#                 # Extract text from shapes
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text") and shape.text.strip():
#                         slide_text.append(shape.text.strip())
#                         slide_shapes.append(shape.text.strip())
                
#                 # Extract notes if configured
#                 if self.powerpoint_config.get("extract_slide_notes", True):
#                     if slide.notes_slide and slide.notes_slide.notes_text_frame:
#                         notes = slide.notes_slide.notes_text_frame.text.strip()
#                         if notes:
#                             slide_text.append(f"[Notes: {notes}]")
                
#                 slides_content.append("\n".join(slide_text))
                
#                 # Create slide summary
#                 slide_title = slide_shapes[0] if slide_shapes else "Untitled"
#                 slide_summaries.append({
#                     "slide_number": slide_idx + 1,
#                     "title": slide_title,
#                     "text_elements": len(slide_shapes)
#                 })
            
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension="pptx",
#                 content_type=["presentation", "powerpoint"],
#                 created_at=datetime.now(),
#                 text_content="\n\n".join(slides_content),
#                 metadata={
#                     "slide_count": len(prs.slides),
#                     "slides": slide_summaries,
#                     "extraction_method": "python-pptx"
#                 }
#             )
            
#         except ImportError:
#             raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
#         except Exception as e:
#             return Document(
#                 doc_id=doc_id,
#                 file_name=file_name,
#                 file_extension="pptx",
#                 content_type=["presentation", "powerpoint", "error"],
#                 created_at=datetime.now(),
#                 text_content=f"Error extracting PowerPoint content: {str(e)}",
#                 metadata={"extraction_error": str(e)}
#             )

# services/processor/src/pipeline/extractors/office_extractor.py
import os
import uuid
from typing import Dict, Any, Optional
from datetime import datetime

from src.models.document import Document

class OfficeExtractor:
    """Extractor for Microsoft Office documents."""
    
    def __init__(self, config: Dict[str, Any] = None):
        """Initialize office extractor."""
        self.config = config or {}
        
    def extract(self, file_path: str) -> Document:
        """Extract content from Office documents."""
        file_extension = os.path.splitext(file_path)[1][1:].lower()
        
        if file_extension in ['xlsx', 'xls']:
            return self._extract_excel(file_path)
        elif file_extension in ['docx', 'doc']:
            return self._extract_word(file_path)
        elif file_extension in ['pptx', 'ppt']:
            return self._extract_powerpoint(file_path)
        else:
            raise ValueError(f"Unsupported Office format: {file_extension}")
    
    def _extract_excel(self, file_path: str) -> Document:
        """Extract content from Excel files."""
        import openpyxl
        import pandas as pd
        
        doc_id = str(uuid.uuid4())
        file_name = os.path.basename(file_path)
        
        document = Document(
            doc_id=doc_id,
            file_name=file_name,
            file_extension="xlsx",
            content_type=["spreadsheet", "excel"],
            created_at=datetime.now(),
            text_content="",
            metadata={}
        )
        
        try:
            # Try pandas first for better data extraction
            excel_file = pd.ExcelFile(file_path)
            
            all_content = []
            sheets_info = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Get sheet info
                sheet_info = {
                    "name": sheet_name,
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": df.columns.tolist()
                }
                sheets_info.append(sheet_info)
                
                # Convert to text
                all_content.append(f"\n=== SHEET: {sheet_name} ===")
                all_content.append(f"Columns: {', '.join(df.columns)}")
                
                # Add data preview (first 10 rows)
                if not df.empty:
                    preview = df.head(10).to_string()
                    all_content.append("\nData Preview:")
                    all_content.append(preview)
                    
                    # Add summary statistics for numeric columns
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    if len(numeric_cols) > 0:
                        all_content.append("\nNumeric Column Summary:")
                        for col in numeric_cols:
                            all_content.append(f"{col}: min={df[col].min()}, max={df[col].max()}, mean={df[col].mean():.2f}")
                
                # Detect data types
                if 'date' in str(df.dtypes).lower() or any('date' in str(col).lower() for col in df.columns):
                    document.content_type.append("time_series")
                
                if len(numeric_cols) > 0:
                    document.content_type.append("numerical_data")
            
            document.text_content = '\n'.join(all_content)
            document.metadata["sheets"] = sheets_info
            document.metadata["total_sheets"] = len(sheets_info)
            
        except Exception as e:
            # Fallback to openpyxl
            try:
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                all_content = []
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    all_content.append(f"\n=== SHEET: {sheet_name} ===")
                    
                    for row in sheet.iter_rows(values_only=True):
                        row_text = []
                        for cell in row:
                            if cell is not None:
                                row_text.append(str(cell))
                        if row_text:
                            all_content.append('\t'.join(row_text))
                
                document.text_content = '\n'.join(all_content)
                
            except Exception as e2:
                document.metadata["extraction_error"] = str(e2)
        
        return document
    
    def _extract_word(self, file_path: str) -> Document:
        """Extract content from Word documents."""
        from docx import Document as DocxDocument
        
        doc_id = str(uuid.uuid4())
        file_name = os.path.basename(file_path)
        
        document = Document(
            doc_id=doc_id,
            file_name=file_name,
            file_extension="docx",
            content_type=["document", "word"],
            created_at=datetime.now(),
            text_content="",
            metadata={}
        )
        
        try:
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            # Extract tables
            tables_text = []
            if doc.tables:
                document.metadata["table_count"] = len(doc.tables)
                for i, table in enumerate(doc.tables):
                    tables_text.append(f"\n[Table {i+1}]")
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(cell.text.strip())
                        tables_text.append(' | '.join(row_text))
            
            # Combine content
            document.text_content = '\n\n'.join(paragraphs)
            if tables_text:
                document.text_content += '\n\n' + '\n'.join(tables_text)
            
            # Extract metadata
            if doc.core_properties.author:
                document.author = doc.core_properties.author
            if doc.core_properties.created:
                document.created_at = doc.core_properties.created
            
            document.metadata["paragraph_count"] = len(paragraphs)
            document.metadata["word_count"] = len(document.text_content.split())
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
        
        return document
    
    def _extract_powerpoint(self, file_path: str) -> Document:
        """Extract content from PowerPoint presentations."""
        from pptx import Presentation
        
        doc_id = str(uuid.uuid4())
        file_name = os.path.basename(file_path)
        
        document = Document(
            doc_id=doc_id,
            file_name=file_name,
            file_extension="pptx",
            content_type=["presentation", "powerpoint"],
            created_at=datetime.now(),
            text_content="",
            metadata={}
        )
        
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            slides_info = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                shapes_count = 0
                has_image = False
                has_chart = False
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    shapes_count += 1
                    
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Check for images
                    if shape.shape_type == 13:  # Picture
                        has_image = True
                    
                    # Check for charts
                    if shape.has_chart:
                        has_chart = True
                        # Extract chart title if available
                        if shape.chart.has_title and shape.chart.title.has_text_frame:
                            slide_text.append(f"[Chart: {shape.chart.title.text_frame.text}]")
                    
                    # Extract table data
                    if shape.has_table:
                        table_text = ["[Table]"]
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            table_text.append(' | '.join(row_text))
                        slide_text.extend(table_text)
                
                # Add slide notes if present
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_text.append(f"\n[Notes: {notes}]")
                
                # Compile slide content
                if slide_text:
                    slides_content.append(f"\n=== SLIDE {i + 1} ===")
                    slides_content.extend(slide_text)
                
                # Store slide metadata
                slides_info.append({
                    "slide_number": i + 1,
                    "shapes_count": shapes_count,
                    "has_image": has_image,
                    "has_chart": has_chart,
                    "text_length": sum(len(t) for t in slide_text)
                })
            
            document.text_content = '\n'.join(slides_content)
            document.metadata["slides"] = slides_info
            document.metadata["total_slides"] = len(prs.slides)
            
            # Detect if it has visual content
            if any(s["has_image"] or s["has_chart"] for s in slides_info):
                document.content_type.append("visual_presentation")
            
        except Exception as e:
            document.metadata["extraction_error"] = str(e)
        
        return document