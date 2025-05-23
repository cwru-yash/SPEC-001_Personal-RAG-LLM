# #!/usr/bin/env python3
# # standalone_dataset_processor.py - Processes your ZIP dataset without complex dependencies

# import os
# import sys
# import zipfile
# import tempfile
# import shutil
# import json
# import uuid
# import logging
# from pathlib import Path
# from datetime import datetime
# from concurrent.futures import ThreadPoolExecutor, as_completed
# from dataclasses import dataclass, asdict
# from typing import Dict, List, Optional, Any

# # Configure logging
# logging.basicConfig(
#     level=logging.INFO,
#     format='%(asctime)s - %(levelname)s - %(message)s',
#     handlers=[
#         logging.StreamHandler(),
#         logging.FileHandler(f'dataset_processing_{datetime.now().strftime("%Y%m%d_%H%M%S")}.log')
#     ]
# )
# logger = logging.getLogger(__name__)

# @dataclass
# class ProcessedDocument:
#     """Simple document representation."""
#     doc_id: str
#     file_name: str
#     category: str
#     doc_type: str  # 'content' or 'metadata'
#     text_content: str
#     metadata: Dict[str, Any]
#     file_size: int
#     processing_time: float
#     extraction_method: str
#     errors: List[str]

# @dataclass
# class DocumentPair:
#     """Represents a pair of content and metadata documents."""
#     pair_id: str
#     base_name: str
#     category: str
#     zip_file: str
#     content_doc: Optional[ProcessedDocument] = None
#     metadata_doc: Optional[ProcessedDocument] = None
#     errors: List[str] = None
    
#     def __post_init__(self):
#         if self.errors is None:
#             self.errors = []

# class SimpleTextExtractor:
#     """Simple text extractor supporting multiple methods."""
    
#     @staticmethod
#     def extract_pdf_text(file_path: str) -> tuple[str, str, List[str]]:
#         """Extract text from PDF using available libraries."""
#         errors = []
        
#         # Try PyMuPDF first
#         try:
#             import fitz
#             with fitz.open(file_path) as doc:
#                 text = ""
#                 for page in doc:
#                     text += page.get_text() + "\n"
#                 return text.strip(), "pymupdf", errors
#         except ImportError:
#             errors.append("PyMuPDF not available")
#         except Exception as e:
#             errors.append(f"PyMuPDF error: {e}")
        
#         # Try pdfplumber
#         try:
#             import pdfplumber
#             with pdfplumber.open(file_path) as pdf:
#                 text = ""
#                 for page in pdf.pages:
#                     page_text = page.extract_text()
#                     if page_text:
#                         text += page_text + "\n"
#                 return text.strip(), "pdfplumber", errors
#         except ImportError:
#             errors.append("pdfplumber not available")
#         except Exception as e:
#             errors.append(f"pdfplumber error: {e}")
        
#         # Fallback: return empty text with errors
#         return "", "none", errors

# class StandaloneDatasetProcessor:
#     """Standalone processor for ZIP datasets."""
    
#     def __init__(self, max_workers: int = 4):
#         """Initialize processor."""
#         self.max_workers = max_workers
#         self.temp_base_dir = tempfile.mkdtemp(prefix="dataset_processing_")
#         self.stats = {
#             "processed_zips": 0,
#             "failed_zips": 0,
#             "document_pairs": 0,
#             "content_docs": 0,
#             "metadata_docs": 0,
#             "by_category": {},
#             "errors": []
#         }
#         logger.info(f"Initialized processor with temp dir: {self.temp_base_dir}")
    
#     def process_dataset(self, input_path: str, output_path: str) -> List[DocumentPair]:
#         """Process the entire dataset."""
        
#         logger.info(f"Starting dataset processing: {input_path} -> {output_path}")
        
#         # Discover ZIP files by category
#         zip_files_by_category = self._discover_zip_files(input_path)
#         total_zips = sum(len(files) for files in zip_files_by_category.values())
        
#         logger.info(f"Found {total_zips} ZIP files across {len(zip_files_by_category)} categories")
        
#         # Process all categories
#         all_pairs = []
#         for category, zip_files in zip_files_by_category.items():
#             logger.info(f"Processing category '{category}' with {len(zip_files)} ZIP files")
#             category_pairs = self._process_category(category, zip_files)
#             all_pairs.extend(category_pairs)
            
#             # Update category stats
#             successful_pairs = len([p for p in category_pairs if p.content_doc and p.metadata_doc])
#             self.stats["by_category"][category] = {
#                 "zip_files": len(zip_files),
#                 "successful_pairs": successful_pairs,
#                 "total_pairs": len(category_pairs)
#             }
        
#         # Save results
#         self._save_results(all_pairs, output_path)
        
#         # Print summary
#         self._print_summary()
        
#         return all_pairs
    
#     def _discover_zip_files(self, input_path: str) -> Dict[str, List[str]]:
#         """Discover ZIP files by category."""
        
#         zip_files_by_category = {}
        
#         for category_dir in os.listdir(input_path):
#             category_path = os.path.join(input_path, category_dir)
            
#             if not os.path.isdir(category_path) or category_dir.startswith('.'):
#                 continue
            
#             zip_files = []
#             for file_name in os.listdir(category_path):
#                 if file_name.endswith('.zip'):
#                     zip_files.append(os.path.join(category_path, file_name))
            
#             if zip_files:
#                 zip_files_by_category[category_dir] = sorted(zip_files)
        
#         return zip_files_by_category
    
#     def _process_category(self, category: str, zip_files: List[str]) -> List[DocumentPair]:
#         """Process all ZIP files in a category."""
        
#         pairs = []
        
#         # Process in parallel
#         with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
#             future_to_zip = {
#                 executor.submit(self._process_single_zip, zip_file, category): zip_file
#                 for zip_file in zip_files
#             }
            
#             for future in as_completed(future_to_zip):
#                 zip_file = future_to_zip[future]
#                 try:
#                     pair = future.result()
#                     pairs.append(pair)
                    
#                     if pair.content_doc and pair.metadata_doc:
#                         self.stats["processed_zips"] += 1
#                         self.stats["document_pairs"] += 1
#                         logger.info(f"✅ Successfully processed: {os.path.basename(zip_file)}")
#                     else:
#                         logger.warning(f"⚠️ Partial processing: {os.path.basename(zip_file)}")
#                         if pair.errors:
#                             logger.warning(f"   Errors: {'; '.join(pair.errors)}")
                            
#                 except Exception as e:
#                     self.stats["failed_zips"] += 1
#                     error_msg = f"Failed to process {zip_file}: {e}"
#                     self.stats["errors"].append(error_msg)
#                     logger.error(f"❌ {error_msg}")
                    
#                     # Create error pair
#                     pairs.append(DocumentPair(
#                         pair_id=str(uuid.uuid4()),
#                         base_name=os.path.splitext(os.path.basename(zip_file))[0],
#                         category=category,
#                         zip_file=zip_file,
#                         errors=[str(e)]
#                     ))
        
#         return pairs
    
#     def _process_single_zip(self, zip_file: str, category: str) -> DocumentPair:
#         """Process a single ZIP file."""
        
#         base_name = os.path.splitext(os.path.basename(zip_file))[0]
#         pair = DocumentPair(
#             pair_id=str(uuid.uuid4()),
#             base_name=base_name,
#             category=category,
#             zip_file=zip_file
#         )
        
#         temp_dir = os.path.join(self.temp_base_dir, f"{base_name}_{uuid.uuid4().hex[:8]}")
        
#         try:
#             # Extract ZIP
#             self._extract_zip(zip_file, temp_dir)
            
#             # Find PDF files (handle nested structure)
#             pdf_files = self._find_pdfs(temp_dir, recursive=True)
            
#             if len(pdf_files) < 2:
#                 pair.errors.append(f"Expected 2 PDFs, found {len(pdf_files)}")
#                 return pair
            
#             # Identify content and metadata PDFs
#             content_pdf, metadata_pdf = self._identify_pdf_pair(pdf_files, base_name)
            
#             # Process both PDFs
#             if content_pdf:
#                 pair.content_doc = self._process_pdf(content_pdf, category, "content", base_name)
#                 if pair.content_doc:
#                     self.stats["content_docs"] += 1
            
#             if metadata_pdf:
#                 pair.metadata_doc = self._process_pdf(metadata_pdf, category, "metadata", base_name)
#                 if pair.metadata_doc:
#                     self.stats["metadata_docs"] += 1
            
#         except Exception as e:
#             pair.errors.append(str(e))
        
#         finally:
#             # Cleanup
#             if os.path.exists(temp_dir):
#                 shutil.rmtree(temp_dir, ignore_errors=True)
        
#         return pair
    
#     def _extract_zip(self, zip_file: str, extract_dir: str):
#         """Extract ZIP file with nested ZIP handling."""
        
#         os.makedirs(extract_dir, exist_ok=True)
        
#         with zipfile.ZipFile(zip_file, 'r') as zip_ref:
#             zip_ref.extractall(extract_dir)
        
#         # Check for nested ZIPs and extract them too
#         for root, _, files in os.walk(extract_dir):
#             for file in files:
#                 if file.endswith('.zip'):
#                     nested_zip = os.path.join(root, file)
#                     nested_extract_dir = os.path.join(root, f"{os.path.splitext(file)[0]}_extracted")
                    
#                     try:
#                         with zipfile.ZipFile(nested_zip, 'r') as nested_zip_ref:
#                             nested_zip_ref.extractall(nested_extract_dir)
#                         # Remove the nested ZIP after extraction
#                         os.remove(nested_zip)
#                     except Exception as e:
#                         logger.warning(f"Failed to extract nested ZIP {nested_zip}: {e}")
    
#     def _find_pdfs(self, directory: str, recursive: bool = True) -> List[str]:
#         """Find all PDF files in directory."""
        
#         pdf_files = []
        
#         if recursive:
#             for root, _, files in os.walk(directory):
#                 for file in files:
#                     if file.lower().endswith('.pdf'):
#                         pdf_files.append(os.path.join(root, file))
#         else:
#             for file in os.listdir(directory):
#                 if file.lower().endswith('.pdf'):
#                     pdf_files.append(os.path.join(directory, file))
        
#         return pdf_files
    
#     def _identify_pdf_pair(self, pdf_files: List[str], base_name: str) -> tuple[Optional[str], Optional[str]]:
#         """Identify content and metadata PDFs."""
        
#         if len(pdf_files) < 2:
#             return None, None
        
#         content_pdf = None
#         metadata_pdf = None
        
#         # Try naming convention first
#         for pdf_file in pdf_files:
#             file_name = os.path.basename(pdf_file)
#             name_without_ext = os.path.splitext(file_name)[0]
            
#             if name_without_ext.endswith('-info') or 'info' in name_without_ext.lower():
#                 metadata_pdf = pdf_file
#             elif name_without_ext == base_name or name_without_ext.startswith(base_name[:6]):
#                 content_pdf = pdf_file
        
#         # Fallback: use file size (metadata usually smaller)
#         if not metadata_pdf or not content_pdf:
#             pdf_sizes = [(f, os.path.getsize(f)) for f in pdf_files]
#             pdf_sizes.sort(key=lambda x: x[1])
            
#             metadata_pdf = pdf_sizes[0][0]  # Smallest
#             content_pdf = pdf_sizes[-1][0]  # Largest
        
#         return content_pdf, metadata_pdf
    
#     def _process_pdf(self, pdf_file: str, category: str, doc_type: str, base_name: str) -> Optional[ProcessedDocument]:
#         """Process a single PDF file."""
        
#         start_time = datetime.now()
        
#         try:
#             # Extract text
#             text_content, extraction_method, errors = SimpleTextExtractor.extract_pdf_text(pdf_file)
            
#             # Get file info
#             file_size = os.path.getsize(pdf_file)
#             processing_time = (datetime.now() - start_time).total_seconds()
            
#             # Create document
#             doc = ProcessedDocument(
#                 doc_id=str(uuid.uuid4()),
#                 file_name=os.path.basename(pdf_file),
#                 category=category,
#                 doc_type=doc_type,
#                 text_content=text_content,
#                 metadata={
#                     "base_name": base_name,
#                     "extraction_method": extraction_method,
#                     "file_size": file_size,
#                     "processed_at": datetime.now().isoformat(),
#                     "text_length": len(text_content),
#                     "word_count": len(text_content.split()) if text_content else 0
#                 },
#                 file_size=file_size,
#                 processing_time=processing_time,
#                 extraction_method=extraction_method,
#                 errors=errors
#             )
            
#             return doc
            
#         except Exception as e:
#             logger.error(f"Error processing PDF {pdf_file}: {e}")
#             return None
    
#     def _save_results(self, pairs: List[DocumentPair], output_path: str):
#         """Save processing results."""
        
#         os.makedirs(output_path, exist_ok=True)
        
#         # Save detailed results
#         results = {
#             "processing_summary": {
#                 "timestamp": datetime.now().isoformat(),
#                 "total_pairs": len(pairs),
#                 "successful_pairs": len([p for p in pairs if p.content_doc and p.metadata_doc]),
#                 "statistics": self.stats
#             },
#             "document_pairs": []
#         }
        
#         for pair in pairs:
#             pair_data = {
#                 "pair_id": pair.pair_id,
#                 "base_name": pair.base_name,
#                 "category": pair.category,
#                 "zip_file": pair.zip_file,
#                 "errors": pair.errors,
#                 "content_document": asdict(pair.content_doc) if pair.content_doc else None,
#                 "metadata_document": asdict(pair.metadata_doc) if pair.metadata_doc else None
#             }
#             results["document_pairs"].append(pair_data)
        
#         # Save main results
#         results_file = os.path.join(output_path, f"processing_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
#         with open(results_file, 'w', encoding='utf-8') as f:
#             json.dump(results, f, indent=2, ensure_ascii=False)
        
#         # Save text content separately for easy access
#         texts_dir = os.path.join(output_path, "extracted_texts")
#         os.makedirs(texts_dir, exist_ok=True)
        
#         for pair in pairs:
#             if pair.content_doc and pair.content_doc.text_content:
#                 content_file = os.path.join(texts_dir, f"{pair.base_name}_content.txt")
#                 with open(content_file, 'w', encoding='utf-8') as f:
#                     f.write(pair.content_doc.text_content)
            
#             if pair.metadata_doc and pair.metadata_doc.text_content:
#                 metadata_file = os.path.join(texts_dir, f"{pair.base_name}_metadata.txt")
#                 with open(metadata_file, 'w', encoding='utf-8') as f:
#                     f.write(pair.metadata_doc.text_content)
        
#         logger.info(f"Results saved to: {results_file}")
#         logger.info(f"Text files saved to: {texts_dir}")
    
#     def _print_summary(self):
#         """Print processing summary."""
        
#         print("\n" + "="*60)
#         print("📊 PROCESSING SUMMARY")
#         print("="*60)
#         print(f"✅ ZIP files processed: {self.stats['processed_zips']}")
#         print(f"❌ ZIP files failed: {self.stats['failed_zips']}")
#         print(f"📄 Document pairs created: {self.stats['document_pairs']}")
#         print(f"📝 Content documents: {self.stats['content_docs']}")
#         print(f"ℹ️  Metadata documents: {self.stats['metadata_docs']}")
        
#         print(f"\n📂 Results by Category:")
#         for category, stats in self.stats['by_category'].items():
#             success_rate = (stats['successful_pairs'] / stats['zip_files']) if stats['zip_files'] > 0 else 0
#             print(f"  {category:15}: {stats['successful_pairs']:3}/{stats['zip_files']:3} ({success_rate:.1%})")
        
#         if self.stats['errors']:
#             print(f"\n⚠️  Errors ({len(self.stats['errors'])}):")
#             for error in self.stats['errors'][:5]:  # Show first 5
#                 print(f"  • {error}")
#             if len(self.stats['errors']) > 5:
#                 print(f"  • ... and {len(self.stats['errors']) - 5} more errors")
    
#     def cleanup(self):
#         """Clean up temporary files."""
#         if os.path.exists(self.temp_base_dir):
#             shutil.rmtree(self.temp_base_dir, ignore_errors=True)
#             logger.info(f"Cleaned up temporary directory: {self.temp_base_dir}")

# def main():
#     """Main processing function."""
    
#     print("🚀 Standalone Dataset Processor")
#     print("="*50)
    
#     # Get paths
#     if len(sys.argv) < 2:
#         print("Usage: python standalone_dataset_processor.py INPUT_PATH [OUTPUT_PATH]")
#         print("\nExample:")
#         print("python standalone_dataset_processor.py /path/to/input /path/to/output")
#         return
    
#     input_path = sys.argv[1]
#     output_path = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(input_path), "processed_output")
    
#     if not os.path.exists(input_path):
#         print(f"❌ Input path does not exist: {input_path}")
#         return
    
#     print(f"📁 Input: {input_path}")
#     print(f"📁 Output: {output_path}")
    
#     # Check dependencies
#     print(f"\n🔍 Checking dependencies...")
#     try:
#         import fitz
#         print("✅ PyMuPDF available")
#     except ImportError:
#         print("⚠️  PyMuPDF not available - install with: pip install pymupdf")
    
#     # Initialize processor
#     processor = StandaloneDatasetProcessor(max_workers=4)
    
#     try:
#         # Process dataset
#         pairs = processor.process_dataset(input_path, output_path)
        
#         print(f"\n🎉 Processing complete!")
#         print(f"📊 Processed {len(pairs)} document pairs")
#         print(f"📁 Results saved to: {output_path}")
        
#     except Exception as e:
#         print(f"❌ Processing failed: {e}")
#         logger.exception("Processing failed")
        
#     finally:
#         # Cleanup
#         processor.cleanup()

# if __name__ == "__main__":
#     main()

#!/usr/bin/env python3
# enhanced_standalone_processor.py - Handles nested ZIP structures for spreadsheets

import os
import sys
import zipfile
import tempfile
import shutil
import json
import uuid
import logging
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, asdict
from typing import Dict, List, Optional, Any

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(f'enhanced_processing_{datetime.now().strftime("%Y%m%d_%H%M%S")}.log')
    ]
)
logger = logging.getLogger(__name__)

@dataclass
class ProcessedDocument:
    """Document representation."""
    doc_id: str
    file_name: str
    category: str
    doc_type: str  # 'content' or 'metadata'
    text_content: str
    metadata: Dict[str, Any]
    file_size: int
    processing_time: float
    extraction_method: str
    errors: List[str]

@dataclass
class DocumentPair:
    """Document pair representation."""
    pair_id: str
    base_name: str
    category: str
    zip_file: str
    content_doc: Optional[ProcessedDocument] = None
    metadata_doc: Optional[ProcessedDocument] = None
    errors: List[str] = None
    
    def __post_init__(self):
        if self.errors is None:
            self.errors = []

class EnhancedTextExtractor:
    """Enhanced text extractor with multiple format support."""
    
    @staticmethod
    def extract_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract text from various file formats."""
        file_ext = os.path.splitext(file_path.lower())[1]
        
        if file_ext == '.pdf':
            return EnhancedTextExtractor.extract_pdf_text(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            return EnhancedTextExtractor.extract_excel_text(file_path)
        elif file_ext in ['.docx', '.doc']:
            return EnhancedTextExtractor.extract_word_text(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return EnhancedTextExtractor.extract_powerpoint_text(file_path)
        elif file_ext in ['.txt', '.csv']:
            return EnhancedTextExtractor.extract_plain_text(file_path)
        else:
            return "", "unsupported", [f"Unsupported file type: {file_ext}"]
    
    @staticmethod
    def extract_pdf_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract text from PDF."""
        errors = []
        
        # Try PyMuPDF first
        try:
            import fitz
            with fitz.open(file_path) as doc:
                text = ""
                for page in doc:
                    text += page.get_text() + "\n"
                return text.strip(), "pymupdf", errors
        except ImportError:
            errors.append("PyMuPDF not available")
        except Exception as e:
            errors.append(f"PyMuPDF error: {e}")
        
        # Try pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text.strip(), "pdfplumber", errors
        except ImportError:
            errors.append("pdfplumber not available")
        except Exception as e:
            errors.append(f"pdfplumber error: {e}")
        
        return "", "none", errors
    
    @staticmethod
    def extract_excel_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract text from Excel files."""
        errors = []
        
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== SHEET: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if any(row_text):  # Skip empty rows
                        text_content.append("\t".join(row_text))
            
            return "\n".join(text_content), "openpyxl", errors
            
        except ImportError:
            errors.append("openpyxl not available - install with: pip install openpyxl")
        except Exception as e:
            errors.append(f"Excel extraction error: {e}")
        
        return "", "none", errors
    
    @staticmethod
    def extract_word_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract text from Word documents."""
        errors = []
        
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            return "\n\n".join(paragraphs), "python-docx", errors
            
        except ImportError:
            errors.append("python-docx not available - install with: pip install python-docx")
        except Exception as e:
            errors.append(f"Word extraction error: {e}")
        
        return "", "none", errors
    
    @staticmethod
    def extract_powerpoint_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract text from PowerPoint presentations."""
        errors = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_text.append(f"=== SLIDE {i + 1} ===\n" + "\n".join(slide_text))
            
            return "\n\n".join(slides_text), "python-pptx", errors
            
        except ImportError:
            errors.append("python-pptx not available - install with: pip install python-pptx")
        except Exception as e:
            errors.append(f"PowerPoint extraction error: {e}")
        
        return "", "none", errors
    
    @staticmethod
    def extract_plain_text(file_path: str) -> tuple[str, str, List[str]]:
        """Extract plain text files."""
        errors = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return content, "plain_text", errors
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='latin-1') as f:
                    content = f.read()
                return content, "plain_text_latin1", errors
            except Exception as e:
                errors.append(f"Text extraction error: {e}")
        except Exception as e:
            errors.append(f"Text extraction error: {e}")
        
        return "", "none", errors

class EnhancedDatasetProcessor:
    """Enhanced processor that handles nested ZIP structures."""
    
    def __init__(self, max_workers: int = 4):
        """Initialize processor."""
        self.max_workers = max_workers
        self.temp_base_dir = tempfile.mkdtemp(prefix="enhanced_processing_")
        self.stats = {
            "processed_zips": 0,
            "failed_zips": 0,
            "document_pairs": 0,
            "content_docs": 0,
            "metadata_docs": 0,
            "by_category": {},
            "errors": []
        }
        logger.info(f"Initialized enhanced processor with temp dir: {self.temp_base_dir}")
    
    def process_dataset(self, input_path: str, output_path: str) -> List[DocumentPair]:
        """Process the entire dataset."""
        
        logger.info(f"Starting enhanced dataset processing: {input_path} -> {output_path}")
        
        # Discover ZIP files by category
        zip_files_by_category = self._discover_zip_files(input_path)
        total_zips = sum(len(files) for files in zip_files_by_category.values())
        
        logger.info(f"Found {total_zips} ZIP files across {len(zip_files_by_category)} categories")
        
        # Process all categories
        all_pairs = []
        for category, zip_files in zip_files_by_category.items():
            logger.info(f"Processing category '{category}' with {len(zip_files)} ZIP files")
            category_pairs = self._process_category(category, zip_files)
            all_pairs.extend(category_pairs)
            
            # Update category stats
            successful_pairs = len([p for p in category_pairs if p.content_doc and p.metadata_doc])
            self.stats["by_category"][category] = {
                "zip_files": len(zip_files),
                "successful_pairs": successful_pairs,
                "total_pairs": len(category_pairs)
            }
        
        # Save results
        self._save_results(all_pairs, output_path)
        
        # Print summary
        self._print_summary()
        
        return all_pairs
    
    def _discover_zip_files(self, input_path: str) -> Dict[str, List[str]]:
        """Discover ZIP files by category."""
        
        zip_files_by_category = {}
        
        for category_dir in os.listdir(input_path):
            category_path = os.path.join(input_path, category_dir)
            
            if not os.path.isdir(category_path) or category_dir.startswith('.'):
                continue
            
            zip_files = []
            for file_name in os.listdir(category_path):
                if file_name.endswith('.zip'):
                    zip_files.append(os.path.join(category_path, file_name))
            
            if zip_files:
                zip_files_by_category[category_dir] = sorted(zip_files)
        
        return zip_files_by_category
    
    def _process_category(self, category: str, zip_files: List[str]) -> List[DocumentPair]:
        """Process all ZIP files in a category."""
        
        pairs = []
        
        # Process in parallel
        with ThreadPoolExecutor(max_workers=self.max_workers) as executor:
            future_to_zip = {
                executor.submit(self._process_single_zip, zip_file, category): zip_file
                for zip_file in zip_files
            }
            
            for future in as_completed(future_to_zip):
                zip_file = future_to_zip[future]
                try:
                    pair = future.result()
                    pairs.append(pair)
                    
                    if pair.content_doc and pair.metadata_doc:
                        self.stats["processed_zips"] += 1
                        self.stats["document_pairs"] += 1
                        logger.info(f"✅ Successfully processed: {os.path.basename(zip_file)}")
                    else:
                        logger.warning(f"⚠️ Partial processing: {os.path.basename(zip_file)}")
                        if pair.errors:
                            logger.warning(f"   Errors: {'; '.join(pair.errors)}")
                            
                except Exception as e:
                    self.stats["failed_zips"] += 1
                    error_msg = f"Failed to process {zip_file}: {e}"
                    self.stats["errors"].append(error_msg)
                    logger.error(f"❌ {error_msg}")
                    
                    # Create error pair
                    pairs.append(DocumentPair(
                        pair_id=str(uuid.uuid4()),
                        base_name=os.path.splitext(os.path.basename(zip_file))[0],
                        category=category,
                        zip_file=zip_file,
                        errors=[str(e)]
                    ))
        
        return pairs
    
    def _process_single_zip(self, zip_file: str, category: str) -> DocumentPair:
        """Process a single ZIP file with enhanced nested handling."""
        
        base_name = os.path.splitext(os.path.basename(zip_file))[0]
        pair = DocumentPair(
            pair_id=str(uuid.uuid4()),
            base_name=base_name,
            category=category,
            zip_file=zip_file
        )
        
        temp_dir = os.path.join(self.temp_base_dir, f"{base_name}_{uuid.uuid4().hex[:8]}")
        
        try:
            # Extract ZIP with enhanced nested handling
            self._extract_zip_enhanced(zip_file, temp_dir, category)
            
            # Find all processable files (PDFs and Office documents)
            processable_files = self._find_processable_files(temp_dir)
            
            if len(processable_files) < 1:
                pair.errors.append(f"No processable files found")
                return pair
            
            # For spreadsheets, handle special case where content is in nested ZIP
            if category == "Spreadsheets":
                content_file, metadata_file = self._identify_spreadsheet_pair(processable_files, base_name)
            else:
                content_file, metadata_file = self._identify_standard_pair(processable_files, base_name)
            
            # Process files
            if content_file:
                pair.content_doc = self._process_file(content_file, category, "content", base_name)
                if pair.content_doc:
                    self.stats["content_docs"] += 1
            
            if metadata_file:
                pair.metadata_doc = self._process_file(metadata_file, category, "metadata", base_name)
                if pair.metadata_doc:
                    self.stats["metadata_docs"] += 1
            
        except Exception as e:
            pair.errors.append(str(e))
            logger.error(f"Error processing {zip_file}: {e}")
        
        finally:
            # Cleanup
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
        
        return pair
    
    def _extract_zip_enhanced(self, zip_file: str, extract_dir: str, category: str):
        """Enhanced ZIP extraction with category-specific handling."""
        
        os.makedirs(extract_dir, exist_ok=True)
        
        with zipfile.ZipFile(zip_file, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)
        
        # For spreadsheets, handle the nested ZIP → Office file structure
        if category == "Spreadsheets":
            self._handle_nested_office_files(extract_dir)
        else:
            # Handle standard nested ZIPs
            self._handle_nested_zips(extract_dir)
    
    def _handle_nested_office_files(self, directory: str):
        """Handle nested office files specifically for spreadsheets."""
        
        for root, _, files in os.walk(directory):
            for file in files:
                if file.endswith('.zip'):
                    nested_zip = os.path.join(root, file)
                    nested_extract_dir = os.path.join(root, f"{os.path.splitext(file)[0]}_content")
                    
                    try:
                        with zipfile.ZipFile(nested_zip, 'r') as nested_zip_ref:
                            nested_zip_ref.extractall(nested_extract_dir)
                        
                        # Check if this extracted office files
                        office_files = []
                        for extracted_root, _, extracted_files in os.walk(nested_extract_dir):
                            for extracted_file in extracted_files:
                                if any(extracted_file.lower().endswith(ext) for ext in ['.xlsx', '.xls', '.docx', '.doc', '.pptx', '.ppt']):
                                    office_files.append(os.path.join(extracted_root, extracted_file))
                        
                        if office_files:
                            # Move office files to main directory for easier access
                            for office_file in office_files:
                                dest_file = os.path.join(root, os.path.basename(office_file))
                                shutil.copy2(office_file, dest_file)
                        
                        # Remove the nested ZIP after extraction
                        os.remove(nested_zip)
                        
                    except Exception as e:
                        logger.warning(f"Failed to extract nested office ZIP {nested_zip}: {e}")
    
    def _handle_nested_zips(self, directory: str):
        """Handle standard nested ZIP files."""
        
        for root, _, files in os.walk(directory):
            for file in files:
                if file.endswith('.zip'):
                    nested_zip = os.path.join(root, file)
                    nested_extract_dir = os.path.join(root, f"{os.path.splitext(file)[0]}_extracted")
                    
                    try:
                        with zipfile.ZipFile(nested_zip, 'r') as nested_zip_ref:
                            nested_zip_ref.extractall(nested_extract_dir)
                        # Remove the nested ZIP after extraction
                        os.remove(nested_zip)
                    except Exception as e:
                        logger.warning(f"Failed to extract nested ZIP {nested_zip}: {e}")
    
    def _find_processable_files(self, directory: str) -> List[str]:
        """Find all processable files (PDFs and Office documents)."""
        
        processable_files = []
        processable_extensions = ['.pdf', '.xlsx', '.xls', '.docx', '.doc', '.pptx', '.ppt', '.txt', '.csv']
        
        for root, _, files in os.walk(directory):
            for file in files:
                if any(file.lower().endswith(ext) for ext in processable_extensions):
                    processable_files.append(os.path.join(root, file))
        
        return processable_files
    
    def _identify_spreadsheet_pair(self, files: List[str], base_name: str) -> tuple[Optional[str], Optional[str]]:
        """Identify content and metadata for spreadsheet category."""
        
        content_file = None
        metadata_file = None
        
        # Look for office files (content) and PDF info files (metadata)
        office_files = [f for f in files if any(f.lower().endswith(ext) for ext in ['.xlsx', '.xls', '.docx', '.doc', '.pptx', '.ppt'])]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        
        # Content is the office file
        if office_files:
            content_file = office_files[0]  # Take the first office file
        
        # Metadata is the PDF with 'info' in name
        for pdf_file in pdf_files:
            if 'info' in os.path.basename(pdf_file).lower():
                metadata_file = pdf_file
                break
        
        # If no info PDF found, look for PDF with base name
        if not metadata_file and pdf_files:
            for pdf_file in pdf_files:
                if base_name in os.path.basename(pdf_file):
                    metadata_file = pdf_file
                    break
        
        return content_file, metadata_file
    
    def _identify_standard_pair(self, files: List[str], base_name: str) -> tuple[Optional[str], Optional[str]]:
        """Identify content and metadata for standard categories."""
        
        if len(files) < 2:
            return files[0] if files else None, None
        
        content_file = None
        metadata_file = None
        
        # Try naming convention first
        for file_path in files:
            file_name = os.path.basename(file_path)
            name_without_ext = os.path.splitext(file_name)[0]
            
            if name_without_ext.endswith('-info') or 'info' in name_without_ext.lower():
                metadata_file = file_path
            elif name_without_ext == base_name or name_without_ext.startswith(base_name[:6]):
                content_file = file_path
        
        # Fallback: use file size (metadata usually smaller)
        if not metadata_file or not content_file:
            file_sizes = [(f, os.path.getsize(f)) for f in files]
            file_sizes.sort(key=lambda x: x[1])
            
            metadata_file = file_sizes[0][0]  # Smallest
            content_file = file_sizes[-1][0]  # Largest
        
        return content_file, metadata_file
    
    def _process_file(self, file_path: str, category: str, doc_type: str, base_name: str) -> Optional[ProcessedDocument]:
        """Process a single file."""
        
        start_time = datetime.now()
        
        try:
            # Extract text using enhanced extractor
            text_content, extraction_method, errors = EnhancedTextExtractor.extract_text(file_path)
            
            # Get file info
            file_size = os.path.getsize(file_path)
            processing_time = (datetime.now() - start_time).total_seconds()
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Create document
            doc = ProcessedDocument(
                doc_id=str(uuid.uuid4()),
                file_name=os.path.basename(file_path),
                category=category,
                doc_type=doc_type,
                text_content=text_content,
                metadata={
                    "base_name": base_name,
                    "extraction_method": extraction_method,
                    "file_size": file_size,
                    "file_extension": file_ext,
                    "processed_at": datetime.now().isoformat(),
                    "text_length": len(text_content),
                    "word_count": len(text_content.split()) if text_content else 0
                },
                file_size=file_size,
                processing_time=processing_time,
                extraction_method=extraction_method,
                errors=errors
            )
            
            return doc
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return None
    
    def _save_results(self, pairs: List[DocumentPair], output_path: str):
        """Save processing results."""
        
        os.makedirs(output_path, exist_ok=True)
        
        # Save detailed results
        results = {
            "processing_summary": {
                "timestamp": datetime.now().isoformat(),
                "total_pairs": len(pairs),
                "successful_pairs": len([p for p in pairs if p.content_doc and p.metadata_doc]),
                "statistics": self.stats
            },
            "document_pairs": []
        }
        
        for pair in pairs:
            pair_data = {
                "pair_id": pair.pair_id,
                "base_name": pair.base_name,
                "category": pair.category,
                "zip_file": pair.zip_file,
                "errors": pair.errors,
                "content_document": asdict(pair.content_doc) if pair.content_doc else None,
                "metadata_document": asdict(pair.metadata_doc) if pair.metadata_doc else None
            }
            results["document_pairs"].append(pair_data)
        
        # Save main results
        results_file = os.path.join(output_path, f"enhanced_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(results, f, indent=2, ensure_ascii=False)
        
        # Save text content separately
        texts_dir = os.path.join(output_path, "extracted_texts")
        os.makedirs(texts_dir, exist_ok=True)
        
        for pair in pairs:
            if pair.content_doc and pair.content_doc.text_content:
                content_file = os.path.join(texts_dir, f"{pair.base_name}_content.txt")
                with open(content_file, 'w', encoding='utf-8') as f:
                    f.write(pair.content_doc.text_content)
            
            if pair.metadata_doc and pair.metadata_doc.text_content:
                metadata_file = os.path.join(texts_dir, f"{pair.base_name}_metadata.txt")
                with open(metadata_file, 'w', encoding='utf-8') as f:
                    f.write(pair.metadata_doc.text_content)
        
        logger.info(f"Enhanced results saved to: {results_file}")
        logger.info(f"Text files saved to: {texts_dir}")
    
    def _print_summary(self):
        """Print processing summary."""
        
        print("\n" + "="*60)
        print("📊 ENHANCED PROCESSING SUMMARY")
        print("="*60)
        print(f"✅ ZIP files processed: {self.stats['processed_zips']}")
        print(f"❌ ZIP files failed: {self.stats['failed_zips']}")
        print(f"📄 Document pairs created: {self.stats['document_pairs']}")
        print(f"📝 Content documents: {self.stats['content_docs']}")
        print(f"ℹ️  Metadata documents: {self.stats['metadata_docs']}")
        
        print(f"\n📂 Results by Category:")
        for category, stats in self.stats['by_category'].items():
            success_rate = (stats['successful_pairs'] / stats['zip_files']) if stats['zip_files'] > 0 else 0
            print(f"  {category:15}: {stats['successful_pairs']:3}/{stats['zip_files']:3} ({success_rate:.1%})")
        
        if self.stats['errors']:
            print(f"\n⚠️  Errors ({len(self.stats['errors'])}):")
            for error in self.stats['errors'][:5]:  # Show first 5
                print(f"  • {error}")
            if len(self.stats['errors']) > 5:
                print(f"  • ... and {len(self.stats['errors']) - 5} more errors")
    
    def cleanup(self):
        """Clean up temporary files."""
        if os.path.exists(self.temp_base_dir):
            shutil.rmtree(self.temp_base_dir, ignore_errors=True)
            logger.info(f"Cleaned up temporary directory: {self.temp_base_dir}")

def main():
    """Main processing function."""
    
    print("🚀 Enhanced Dataset Processor (Handles Nested Office Files)")
    print("="*60)
    
    # Get paths
    if len(sys.argv) < 2:
        print("Usage: python enhanced_standalone_processor.py INPUT_PATH [OUTPUT_PATH]")
        print("\nExample:")
        print("python enhanced_standalone_processor.py /path/to/input /path/to/output")
        return
    
    input_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(input_path), "enhanced_output")
    
    if not os.path.exists(input_path):
        print(f"❌ Input path does not exist: {input_path}")
        return
    
    print(f"📁 Input: {input_path}")
    print(f"📁 Output: {output_path}")
    
    # Check dependencies
    print(f"\n🔍 Checking dependencies...")
    deps = [
        ("PyMuPDF", "fitz", "pip install pymupdf"),
        ("openpyxl", "openpyxl", "pip install openpyxl"),
        ("python-docx", "docx", "pip install python-docx"),
        ("python-pptx", "pptx", "pip install python-pptx")
    ]
    
    for name, module, install_cmd in deps:
        try:
            __import__(module)
            print(f"✅ {name} available")
        except ImportError:
            print(f"⚠️  {name} not available - install with: {install_cmd}")
    
    # Initialize processor
    processor = EnhancedDatasetProcessor(max_workers=4)
    
    try:
        # Process dataset
        pairs = processor.process_dataset(input_path, output_path)
        
        print(f"\n🎉 Enhanced processing complete!")
        print(f"📊 Processed {len(pairs)} document pairs")
        print(f"📁 Results saved to: {output_path}")
        
    except Exception as e:
        print(f"❌ Processing failed: {e}")
        logger.exception("Processing failed")
        
    finally:
        # Cleanup
        processor.cleanup()

if __name__ == "__main__":
    main()